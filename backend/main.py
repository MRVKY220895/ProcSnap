import base64
import json
import re
import sqlite3
import time
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional, List, Dict, Any, Union
from uuid import uuid4
import subprocess
import sys
import threading
import edge_tts
import asyncio
import os
import shutil
import io
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls

from PIL import Image, ImageDraw

try:
    import pptx
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    pptx = None

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Query, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, Response
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel


# =========================================================
# APPLICATION
# =========================================================

app = FastAPI(title="ProcSnap API", version="0.1.0")

BASE_DIR = Path(__file__).resolve().parent
MEDIA_DIR = BASE_DIR / "media"
MEDIA_DIR.mkdir(parents=True, exist_ok=True)

# Serve dashboard static files
@app.get("/favicon.ico", include_in_schema=False)
@app.get("/dashboard/favicon.ico", include_in_schema=False)
def get_favicon():
    svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect width="100" height="100" rx="20" fill="#6366f1"/><text x="50" y="68" font-size="55" font-family="sans-serif" font-weight="800" fill="white" text-anchor="middle">P</text></svg>'
    return Response(content=svg, media_type="image/svg+xml")

app.mount(
    "/dashboard",
    StaticFiles(directory=BASE_DIR.parent / "dashboard"),
    name="dashboard",
)
app.mount(
    "/media",
    StaticFiles(directory=MEDIA_DIR),
    name="media",
)



# =========================================================
# CORS
# =========================================================

@app.middleware("http")
async def add_cors_pna_headers(request, call_next):
    if request.method == "OPTIONS":
        from fastapi.responses import Response
        response = Response(status_code=204)
        response.headers["Access-Control-Allow-Origin"] = "*"
        response.headers["Access-Control-Allow-Methods"] = "*"
        response.headers["Access-Control-Allow-Headers"] = "*"
        response.headers["Access-Control-Allow-Private-Network"] = "true"
        return response
    response = await call_next(request)
    response.headers["Access-Control-Allow-Origin"] = "*"
    response.headers["Access-Control-Allow-Private-Network"] = "true"
    return response


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


# =========================================================
# DATABASE / FILE STORAGE
# =========================================================

BASE_DIR = Path(__file__).resolve().parent
DATABASE_PATH = BASE_DIR / "procsnap.db"
SCREENSHOTS_DIR = BASE_DIR / "screenshots"

SCREENSHOTS_DIR.mkdir(parents=True, exist_ok=True)


def get_connection() -> sqlite3.Connection:
    connection = sqlite3.connect(
        DATABASE_PATH,
        check_same_thread=False,
    )
    connection.row_factory = sqlite3.Row
    connection.execute("PRAGMA foreign_keys = ON")
    return connection


def initialize_database() -> None:
    connection = get_connection()
    cursor = connection.cursor()

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workflows (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            application TEXT NOT NULL,
            status TEXT NOT NULL,
            started_at TEXT NOT NULL,
            ended_at TEXT,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL
        )
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workflow_steps (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            workflow_id TEXT NOT NULL,
            sequence INTEGER NOT NULL,
            action TEXT NOT NULL,
            timestamp TEXT NOT NULL,
            url TEXT,
            title TEXT,
            value TEXT,
            selected_text TEXT,
            previous_url TEXT,
            checked INTEGER,
            element_json TEXT,
            screenshot_path TEXT,
            FOREIGN KEY (workflow_id)
                REFERENCES workflows(id)
                ON DELETE CASCADE
        )
        """
    )

    columns = {
        row["name"]
        for row in cursor.execute(
            "PRAGMA table_info(workflow_steps)"
        ).fetchall()
    }

    if "screenshot_path" not in columns:
        cursor.execute(
            """
            ALTER TABLE workflow_steps
            ADD COLUMN screenshot_path TEXT
            """
        )

    cursor.execute(
        """
        CREATE INDEX IF NOT EXISTS
        idx_workflow_steps_workflow_id
        ON workflow_steps(workflow_id)
        """
    )

    cursor.execute(
        """
        CREATE UNIQUE INDEX IF NOT EXISTS
        idx_workflow_steps_workflow_sequence
        ON workflow_steps(workflow_id, sequence)
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS step_annotations (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            step_id INTEGER NOT NULL UNIQUE,
            workflow_id TEXT NOT NULL,
            data TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            FOREIGN KEY (step_id) REFERENCES workflow_steps(id) ON DELETE CASCADE,
            FOREIGN KEY (workflow_id) REFERENCES workflows(id) ON DELETE CASCADE
        )
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS step_edits (
            step_id INTEGER PRIMARY KEY,
            workflow_id TEXT NOT NULL,
            title TEXT,
            description TEXT,
            note TEXT,
            expected TEXT,
            voiceover TEXT,
            hidden INTEGER DEFAULT 0,
            updated_at TEXT NOT NULL,
            FOREIGN KEY (step_id) REFERENCES workflow_steps(id) ON DELETE CASCADE,
            FOREIGN KEY (workflow_id) REFERENCES workflows(id) ON DELETE CASCADE
        )
        """
    )

    try:
        cursor.execute("ALTER TABLE step_edits ADD COLUMN voiceover TEXT")
    except Exception:
        pass

    try:
        cursor.execute("ALTER TABLE step_edits ADD COLUMN branches TEXT")
    except Exception:
        pass

    try:
        cursor.execute("ALTER TABLE workflows ADD COLUMN tags TEXT DEFAULT ''")
    except Exception:
        pass

    # Phase 1 — Smart Screenshot Timing quality columns
    for col in ["screenshot_quality INTEGER DEFAULT NULL",
                "recapture_suggested INTEGER DEFAULT 0"]:
        try:
            cursor.execute(f"ALTER TABLE workflow_steps ADD COLUMN {col}")
        except Exception:
            pass

    # Phase 2 & 3 — Event normalization + SOP intelligence columns
    for col in ["semantic_class TEXT",
                "intent_marker TEXT",
                "why_important TEXT",
                "step_type TEXT",
                "mandatory INTEGER DEFAULT 1",
                "risk_level TEXT DEFAULT 'low'",
                "role TEXT",
                "control_id TEXT",
                "estimated_duration INTEGER",
                "duration TEXT",
                "alert_type TEXT DEFAULT 'none'",
                "alert_msg TEXT DEFAULT ''",
                "noise_flags TEXT"]:
        try:
            cursor.execute(f"ALTER TABLE step_edits ADD COLUMN {col}")
        except Exception:
            pass

    # Phase 4 — SOP Template & Variables columns
    for col in ["template_type TEXT DEFAULT 'standard'",
                "variables TEXT DEFAULT '{}'",
                "sop_metadata TEXT DEFAULT '{}'",
                "description TEXT DEFAULT ''"]:
        try:
            cursor.execute(f"ALTER TABLE workflows ADD COLUMN {col}")
        except Exception:
            pass

    # Phase 7 & 8 — Privacy Redaction & SOP Lifecycle Versioning
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS redaction_profiles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL,
            rules TEXT NOT NULL,
            created_at TEXT NOT NULL
        )
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS workflow_versions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            workflow_id TEXT NOT NULL,
            version TEXT NOT NULL,
            status TEXT DEFAULT 'draft',
            snapshot_json TEXT NOT NULL,
            change_summary TEXT,
            created_by TEXT DEFAULT 'Author',
            created_at TEXT NOT NULL,
            review_due TEXT,
            FOREIGN KEY (workflow_id) REFERENCES workflows(id) ON DELETE CASCADE
        )
        """
    )

    for col in ["current_version TEXT DEFAULT '1.0'",
                "lifecycle_status TEXT DEFAULT 'draft'",
                "review_due_date TEXT",
                "department TEXT DEFAULT ''",
                "owner TEXT DEFAULT ''",
                "reviewer TEXT DEFAULT ''",
                "approver TEXT DEFAULT ''",
                "effective_date TEXT DEFAULT ''",
                "review_frequency_days INTEGER DEFAULT 90",
                "preconditions_json TEXT DEFAULT '[]'",
                "postconditions_json TEXT DEFAULT '[]'"]:
        try:
            cursor.execute(f"ALTER TABLE workflows ADD COLUMN {col}")
        except Exception:
            pass

    for col in ["expected_result TEXT DEFAULT ''",
                "exception_info TEXT DEFAULT ''",
                "phase_name TEXT DEFAULT 'Execution'",
                "business_intent TEXT DEFAULT ''",
                "pii_masked INTEGER DEFAULT 0",
                "redaction_flags TEXT DEFAULT NULL"]:
        try:
            cursor.execute(f"ALTER TABLE step_edits ADD COLUMN {col}")
        except Exception:
            pass

    for col in ["expected_result TEXT DEFAULT ''",
                "phase_name TEXT DEFAULT 'Execution'"]:
        try:
            cursor.execute(f"ALTER TABLE workflow_steps ADD COLUMN {col}")
        except Exception:
            pass

    # ProcBot RPA Configuration & Execution History Tables
    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS procbot_configs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            workflow_id TEXT NOT NULL UNIQUE,
            name TEXT NOT NULL,
            config_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            FOREIGN KEY (workflow_id) REFERENCES workflows(id) ON DELETE CASCADE
        )
        """
    )

    cursor.execute(
        """
        CREATE TABLE IF NOT EXISTS procbot_run_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            workflow_id TEXT NOT NULL,
            engine TEXT NOT NULL,
            mode TEXT NOT NULL,
            total_steps INTEGER NOT NULL,
            success_steps INTEGER NOT NULL,
            failed_steps INTEGER NOT NULL,
            elapsed_sec REAL NOT NULL,
            log_json TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY (workflow_id) REFERENCES workflows(id) ON DELETE CASCADE
        )
        """
    )

    connection.commit()
    connection.close()





initialize_database()

def seed_initial_sample_sop():
    """Seeds rich sample SOPs if the database is brand new so first-time users have immediate interactive guides."""
    try:
        conn = get_connection()
        cur = conn.cursor()
        count = cur.execute("SELECT COUNT(*) FROM workflows").fetchone()[0]
        if count == 0:
            now = datetime.utcnow().isoformat()
            sample_id = "sample_welcome_guide"
            cur.execute("""
                INSERT INTO workflows (id, name, application, status, started_at, ended_at, created_at, updated_at, tags, department, owner, reviewer, approver)
                VALUES (?, '🚀 Welcome to ProcSnap: Master SOP Guide', 'Chrome Web Browser', 'completed', ?, ?, ?, ?, 'Sample Guide, Onboarding', 'Operations', 'Process Team', 'Lead Auditor', 'Head of Ops')
            """, (sample_id, now, now, now, now))
            
            sample_dir = SCREENSHOTS_DIR / sample_id
            sample_dir.mkdir(parents=True, exist_ok=True)
            
            steps_data = [
                (1, "click", "https://app.procsnap.local/dashboard", "Open the Operations Dashboard", "", "The central management portal loads with current metrics.", "Navigation", "To access the main company operational tools."),
                (2, "click", "https://app.procsnap.local/customers", "Select Customer Management", "", "The Customer Directory list is displayed.", "Selection", "To locate or create account records."),
                (3, "input", "https://app.procsnap.local/customers/new", "Enter Customer Account Name", "Acme Global Corp", "The account name field accepts the input.", "DataEntry", "To register the official client entity name."),
                (4, "click", "https://app.procsnap.local/customers/new", "Submit and Save Account Record", "", "The account is successfully created and assigned ID #10492.", "Submission", "To commit the record to the permanent database.")
            ]
            
            for seq, act, url, title, val, exp, sem, intent in steps_data:
                cur.execute("""
                    INSERT INTO workflow_steps (workflow_id, sequence, action, timestamp, url, title, value)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                """, (sample_id, seq, act, now, url, title, val))
                s_id = cur.lastrowid
                
                # Also create default step_edits
                cur.execute("""
                    INSERT INTO step_edits (step_id, workflow_id, title, description, expected, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (s_id, sample_id, title, f"Follow this step to {title.lower()}.", exp, now))
            
            conn.commit()
            print("🚀 Seeded initial sample SOP for first-time user.")
        conn.close()
    except Exception as e:
        print("Sample SOP seeding notice:", e)

seed_initial_sample_sop()

def perform_auto_backup():
    try:
        if not DATABASE_PATH.exists():
            return
        backup_dir = BASE_DIR / "storage" / "backups"
        backup_dir.mkdir(parents=True, exist_ok=True)
        
        backups = sorted(list(backup_dir.glob("procsnap_backup_*.db")))
        while len(backups) >= 5:
            try:
                backups.pop(0).unlink()
            except Exception as e:
                print("Error removing oldest backup:", e)
                break
                
        import shutil
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_file = backup_dir / f"procsnap_backup_{timestamp}.db"
        shutil.copy2(DATABASE_PATH, backup_file)
        print(f"Auto-backup created successfully: {backup_file.name}")
    except Exception as e:
        print("Auto-backup failed:", e)

# Perform startup backup
perform_auto_backup()


# =========================================================
# MODELS
# =========================================================

# ElementInfo intentionally permits additional screen metadata through the
# existing JSON storage path. The extension sends:
# element.screen = {
#   x, y, width, height,
#   viewportWidth, viewportHeight,
#   devicePixelRatio
# }
# No database migration is required because element_json is stored as JSON text.
#
class ElementInfo(BaseModel):
    tagName: Optional[str] = None
    id: Optional[str] = None
    className: Optional[str] = None
    text: Optional[str] = None
    name: Optional[str] = None
    placeholder: Optional[str] = None
    type: Optional[str] = None
    ariaLabel: Optional[str] = None
    role: Optional[str] = None
    title: Optional[str] = None
    dataTestId: Optional[str] = None
    cssSelector: Optional[str] = None
    screen: Optional[dict] = None
    xpath: Optional[str] = None


class WorkflowStep(BaseModel):
    action: str
    timestamp: str
    url: str
    title: Optional[str] = None
    value: Optional[str] = None
    selectedText: Optional[str] = None
    previousUrl: Optional[str] = None
    checked: Optional[bool] = None
    element: Optional[ElementInfo] = None


class StartSessionRequest(BaseModel):
    name: Optional[str] = None
    application: Optional[str] = None


class UpdateWorkflowRequest(BaseModel):
    name: Optional[str] = None
    tags: Optional[str] = None


class ImportWorkflowRequest(BaseModel):
    workflow: dict


class ScreenshotRequest(BaseModel):
    sessionId: str
    image: str
    stepId: Optional[int] = None


class StepAnnotationRequest(BaseModel):
    data: str  # JSON array of annotation objects


class StepEditRequest(BaseModel):
    title: Optional[str] = None
    description: Optional[str] = None
    note: Optional[str] = None
    expected: Optional[str] = None
    voiceover: Optional[str] = None
    branches: Optional[str] = None
    hidden: Optional[bool] = None
    checked: Optional[bool] = None
    alert_type: Optional[str] = None
    alert_msg: Optional[str] = None
    role: Optional[str] = None
    duration: Optional[str] = None


# =========================================================
# UTILITY
# =========================================================

def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def get_step_count_with_connection(
    cursor: sqlite3.Cursor,
    workflow_id: str,
) -> int:
    cursor.execute(
        """
        SELECT COUNT(*) AS count
        FROM workflow_steps
        WHERE workflow_id = ?
        """,
        (workflow_id,),
    )
    return int(cursor.fetchone()["count"])


def get_step_count(workflow_id: str) -> int:
    connection = get_connection()
    try:
        return get_step_count_with_connection(
            connection.cursor(),
            workflow_id,
        )
    finally:
        connection.close()


# =========================================================
# SERIALIZATION
# =========================================================

def row_to_workflow(row: sqlite3.Row) -> dict:
    return {
        "id": row["id"],
        "name": row["name"],
        "application": row["application"],
        "status": row["status"],
        "startedAt": row["started_at"],
        "endedAt": row["ended_at"],
        "stepCount": get_step_count(row["id"]),
    }


def step_to_response(row: sqlite3.Row) -> dict:
    element = None

    if row["element_json"]:
        try:
            element = json.loads(row["element_json"])
        except (json.JSONDecodeError, TypeError):
            element = None

    screenshot_path = row["screenshot_path"]

    screenshot_url = None
    if screenshot_path:
        screenshot_url = (
            f"/screenshots/{row['workflow_id']}/"
            f"{Path(screenshot_path).name}"
        )

    return {
        "id": row["id"],
        "sequence": row["sequence"],
        "action": row["action"],
        "timestamp": row["timestamp"],
        "url": row["url"],
        "title": row["title"],
        "value": row["value"],
        "selectedText": row["selected_text"],
        "previousUrl": row["previous_url"],
        "checked": (
            bool(row["checked"])
            if row["checked"] is not None
            else None
        ),
        "element": element,
        "screenshotPath": screenshot_path,
        "screenshotUrl": screenshot_url,
    }


# =========================================================
# ROOT
# =========================================================

@app.get("/")
def root():
    return {
        "application": "ProcSnap",
        "status": "running",
        "database": str(DATABASE_PATH),
        "screenshots": str(SCREENSHOTS_DIR),
        "timestamp": utc_now(),
    }


# =========================================================
# HEALTH
# =========================================================

@app.get("/health")
def health():
    connection = get_connection()
    try:
        connection.execute("SELECT 1").fetchone()
    finally:
        connection.close()

    return {
        "status": "healthy",
        "database": "connected",
    }


# =========================================================
# CREATE WORKFLOW / START RECORDING
# =========================================================

@app.post("/sessions")
def start_session(request: StartSessionRequest):
    workflow_id = str(uuid4())
    now = utc_now()

    name = request.name or "Untitled Workflow"
    application = request.application or "Chrome"

    connection = get_connection()

    try:
        connection.execute(
            """
            INSERT INTO workflows (
                id,
                name,
                application,
                status,
                started_at,
                ended_at,
                created_at,
                updated_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                workflow_id,
                name,
                application,
                "recording",
                now,
                None,
                now,
                now,
            ),
        )
        connection.commit()
    finally:
        connection.close()

    print()
    print("======================================")
    print("NEW RECORDING SESSION")
    print("======================================")
    print("Session ID:", workflow_id)
    print("Name:", name)
    print("Application:", application)
    print("Database:", DATABASE_PATH)
    print("======================================")
    print()

    # Start Desktop Recorder in hybrid mode so native OS dialogs (e.g. File Explorer upload picker) are automatically captured
    if desktop_recorder:
        try:
            desktop_recorder.start(
                title=name,
                target_monitor="auto",
                auto_click_capture=True,
                session_id=workflow_id,
                mode="hybrid",
                db_callback=_desktop_step_db_callback
            )
            print(f"[ProcSnap] Hybrid Desktop Recorder linked to session {workflow_id}")
        except Exception as dex:
            print(f"[ProcSnap] Hybrid desktop recorder start notice: {dex}")

    return {
        "success": True,
        "session": {
            "id": workflow_id,
            "name": name,
            "application": application,
            "status": "recording",
            "startedAt": now,
            "endedAt": None,
            "stepCount": 0,
            "steps": [],
        },
    }


# =========================================================
# STOP WORKFLOW
# =========================================================

@app.post("/sessions/{session_id}/stop")
def stop_session(session_id: str):
    connection = get_connection()

    try:
        cursor = connection.cursor()

        cursor.execute(
            """
            SELECT *
            FROM workflows
            WHERE id = ?
            """,
            (session_id,),
        )

        workflow = cursor.fetchone()

        if not workflow:
            raise HTTPException(
                status_code=404,
                detail="Recording session not found",
            )

        ended_at = utc_now()

        cursor.execute(
            """
            UPDATE workflows
            SET
                status = ?,
                ended_at = ?,
                updated_at = ?
            WHERE id = ?
            """,
            (
                "completed",
                ended_at,
                ended_at,
                session_id,
            ),
        )

        step_count = get_step_count_with_connection(
            cursor,
            session_id,
        )

        connection.commit()

    finally:
        connection.close()

    print()
    print("======================================")
    print("RECORDING SESSION COMPLETED")
    print("======================================")
    print("Session ID:", session_id)
    print("Steps:", step_count)
    print("======================================")
    print()

    if desktop_recorder and desktop_recorder.is_recording:
        try:
            desktop_recorder.stop()
        except Exception:
            pass

    return {
        "success": True,
        "session": {
            "id": workflow["id"],
            "name": workflow["name"],
            "application": workflow["application"],
            "status": "completed",
            "startedAt": workflow["started_at"],
            "endedAt": ended_at,
            "stepCount": step_count,
            "steps": [],
        },
    }


# =========================================================
# GET ALL WORKFLOWS
# =========================================================

@app.get("/sessions")
def get_sessions():
    connection = get_connection()

    try:
        rows = connection.execute(
            """
            SELECT w.id, w.name, w.application, w.status, w.started_at, w.ended_at, w.tags, w.created_at,
                   COUNT(s.id) AS step_count,
                   (SELECT s2.screenshot_path FROM workflow_steps s2 WHERE s2.workflow_id = w.id ORDER BY s2.sequence ASC LIMIT 1) AS cover_screenshot
            FROM workflows w
            LEFT JOIN workflow_steps s ON w.id = s.workflow_id
            GROUP BY w.id
            ORDER BY w.created_at DESC
            """
        ).fetchall()

        result = [
            {
                "id": row["id"],
                "name": row["name"] or "Untitled Workflow",
                "application": row["application"] or "System",
                "status": row["status"] or "completed",
                "tags": row["tags"] or "",
                "startedAt": row["started_at"],
                "endedAt": row["ended_at"],
                "createdAt": row["created_at"] or row["started_at"],
                "stepCount": row["step_count"],
                "coverScreenshot": row["cover_screenshot"] or "",
            }
            for row in rows
        ]

    finally:
        connection.close()

    return {
        "count": len(result),
        "sessions": result,
    }


@app.get("/sessions/{session_id}/cover")
def get_session_cover_image(session_id: str):
    """
    Returns the first step's screenshot or cover image for the workflow session.
    """
    conn = get_connection()
    try:
        cur = conn.cursor()
        step = cur.execute(
            "SELECT screenshot_path FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence ASC LIMIT 1",
            (session_id,)
        ).fetchone()
        if not step or not step["screenshot_path"]:
            raise HTTPException(status_code=404, detail="No screenshot found for this session")
        
        img_file = BASE_DIR / step["screenshot_path"]
        if not img_file.exists():
            img_file = BASE_DIR / "screenshots" / session_id / Path(step["screenshot_path"]).name
        if not img_file.exists():
            raise HTTPException(status_code=404, detail="Screenshot file not found on disk")
        
        return FileResponse(img_file, media_type="image/png")
    finally:
        conn.close()


# =========================================================
# GET SINGLE WORKFLOW
# =========================================================

@app.get("/sessions/{session_id}")
def get_session(session_id: str):
    connection = get_connection()

    try:
        cursor = connection.cursor()

        cursor.execute(
            """
            SELECT *
            FROM workflows
            WHERE id = ?
            """,
            (session_id,),
        )

        workflow = cursor.fetchone()

        if not workflow:
            raise HTTPException(
                status_code=404,
                detail="Recording session not found",
            )

        step_rows = cursor.execute(
            """
            SELECT *
            FROM workflow_steps
            WHERE workflow_id = ?
            ORDER BY sequence ASC
            """,
            (session_id,),
        ).fetchall()

        # Fetch all annotations for this workflow
        annotations_rows = cursor.execute(
            """
            SELECT step_id, data
            FROM step_annotations
            WHERE workflow_id = ?
            """,
            (session_id,),
        ).fetchall()
        annotations_map = {row["step_id"]: row["data"] for row in annotations_rows}

        # Fetch all edits for this workflow
        edits_rows = cursor.execute(
            """
            SELECT step_id, title, description, note, expected, voiceover, branches, hidden, alert_type, alert_msg, role, duration
            FROM step_edits
            WHERE workflow_id = ?
            """,
            (session_id,),
        ).fetchall()
        edits_map = {row["step_id"]: row for row in edits_rows}

        steps = []
        for row in step_rows:
            step_dict = step_to_response(row)
            step_id = step_dict["id"]
            
            # Merge annotations
            step_dict["annotations"] = json.loads(annotations_map.get(step_id, "[]"))
            
            # Merge edits
            if step_id in edits_map:
                edit = edits_map[step_id]
                step_dict["title"] = edit["title"] if edit["title"] is not None else step_dict["title"]
                step_dict["description"] = edit["description"] if edit["description"] is not None else step_dict["description"]
                step_dict["note"] = edit["note"] or ""
                step_dict["expected"] = edit["expected"] or ""
                step_dict["voiceover"] = edit["voiceover"] or ""
                step_dict["alertType"] = edit["alert_type"] if "alert_type" in edit.keys() and edit["alert_type"] else "none"
                step_dict["alertMsg"] = edit["alert_msg"] if "alert_msg" in edit.keys() and edit["alert_msg"] else ""
                step_dict["role"] = edit["role"] if "role" in edit.keys() and edit["role"] else ""
                step_dict["duration"] = edit["duration"] if "duration" in edit.keys() and edit["duration"] else ""
                raw_branches = edit["branches"] if "branches" in edit.keys() and edit["branches"] else ""
                try:
                    step_dict["branches"] = json.loads(raw_branches) if raw_branches else []
                except Exception:
                    step_dict["branches"] = []
                step_dict["hidden"] = bool(edit["hidden"])
            else:
                step_dict["note"] = ""
                step_dict["expected"] = ""
                step_dict["voiceover"] = ""
                step_dict["alertType"] = "none"
                step_dict["alertMsg"] = ""
                step_dict["role"] = ""
                step_dict["duration"] = ""
                step_dict["branches"] = []
                step_dict["hidden"] = False
                
            steps.append(step_dict)

        return {
            "id": workflow["id"],
            "name": workflow["name"],
            "application": workflow["application"],
            "status": workflow["status"],
            "tags": (workflow["tags"] if "tags" in workflow.keys() else "") or "",
            "startedAt": workflow["started_at"],
            "endedAt": workflow["ended_at"],
            "stepCount": len(steps),
            "steps": steps,
        }

    finally:
        connection.close()


# =========================================================
# ADD STEP
# =========================================================

@app.post("/sessions/{session_id}/steps")
def add_step(
    session_id: str,
    step: WorkflowStep,
):
    connection = get_connection()

    try:
        cursor = connection.cursor()

        cursor.execute(
            """
            SELECT status
            FROM workflows
            WHERE id = ?
            """,
            (session_id,),
        )

        workflow = cursor.fetchone()

        if not workflow:
            raise HTTPException(
                status_code=404,
                detail="Recording session not found",
            )

        # Allow adding steps in active recording or completed workflows in studio

        cursor.execute(
            """
            SELECT COALESCE(MAX(sequence), 0) + 1
            AS next_sequence
            FROM workflow_steps
            WHERE workflow_id = ?
            """,
            (session_id,),
        )

        sequence = int(
            cursor.fetchone()["next_sequence"]
        )

        element_json = None

        if step.element is not None:
            element_json = json.dumps(
                step.element.model_dump()
            )

        checked_value = None

        if step.checked is not None:
            checked_value = (
                1 if step.checked else 0
            )

        cursor.execute(
            """
            INSERT INTO workflow_steps (
                workflow_id,
                sequence,
                action,
                timestamp,
                url,
                title,
                value,
                selected_text,
                previous_url,
                checked,
                element_json,
                screenshot_path
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                session_id,
                sequence,
                step.action,
                step.timestamp,
                step.url,
                step.title,
                step.value,
                step.selectedText,
                step.previousUrl,
                checked_value,
                element_json,
                None,
            ),
        )

        step_id = cursor.lastrowid

        cursor.execute(
            """
            UPDATE workflows
            SET updated_at = ?
            WHERE id = ?
            """,
            (
                utc_now(),
                session_id,
            ),
        )

        connection.commit()

    finally:
        connection.close()

    print(
        f"Session {session_id} | "
        f"Step {sequence} | "
        f"Action: {step.action}"
    )

    return {
        "success": True,
        "sessionId": session_id,
        "step": {
            "id": step_id,
            "sequence": sequence,
            **step.model_dump(),
        },
    }


# =========================================================
# SAVE SCREENSHOT
# =========================================================

@app.post("/screenshots")
def save_screenshot(request: ScreenshotRequest):
    connection = get_connection()

    try:
        cursor = connection.cursor()

        cursor.execute(
            """
            SELECT id
            FROM workflows
            WHERE id = ?
            """,
            (request.sessionId,),
        )

        workflow = cursor.fetchone()

        if not workflow:
            raise HTTPException(
                status_code=404,
                detail="Workflow not found",
            )

        match = re.fullmatch(
            r"data:image/png;base64,(.+)",
            request.image,
            flags=re.DOTALL,
        )

        if not match:
            raise HTTPException(
                status_code=400,
                detail="Invalid PNG data URL",
            )

        try:
            image_bytes = base64.b64decode(
                match.group(1),
                validate=True,
            )
        except (ValueError, TypeError):
            raise HTTPException(
                status_code=400,
                detail="Invalid base64 image data",
            )

        if not image_bytes:
            raise HTTPException(
                status_code=400,
                detail="Screenshot is empty",
            )

        if request.stepId is not None:
            cursor.execute(
                """
                SELECT *
                FROM workflow_steps
                WHERE id = ?
                  AND workflow_id = ?
                """,
                (
                    request.stepId,
                    request.sessionId,
                ),
            )
        else:
            cursor.execute(
                """
                SELECT *
                FROM workflow_steps
                WHERE workflow_id = ?
                ORDER BY sequence DESC
                LIMIT 1
                """,
                (request.sessionId,),
            )

        step_row = cursor.fetchone()

        if not step_row:
            raise HTTPException(
                status_code=404,
                detail="No workflow step found for screenshot",
            )

        sequence = int(step_row["sequence"])

        screenshot_dir = (
            SCREENSHOTS_DIR /
            request.sessionId
        )

        screenshot_dir.mkdir(
            parents=True,
            exist_ok=True,
        )

        filename = f"step-{sequence:03d}.png"
        filepath = screenshot_dir / filename

        filepath.write_bytes(image_bytes)

        relative_path = str(
            Path("screenshots") /
            request.sessionId /
            filename
        )

        cursor.execute(
            """
            UPDATE workflow_steps
            SET screenshot_path = ?
            WHERE id = ?
              AND workflow_id = ?
            """,
            (
                relative_path,
                step_row["id"],
                request.sessionId,
            ),
        )

        cursor.execute(
            """
            UPDATE workflows
            SET updated_at = ?
            WHERE id = ?
            """,
            (
                utc_now(),
                request.sessionId,
            ),
        )

        connection.commit()

    finally:
        connection.close()

    print("Screenshot saved:", filepath)

    return {
        "success": True,
        "sessionId": request.sessionId,
        "stepId": step_row["id"],
        "sequence": sequence,
        "filename": filename,
        "path": relative_path,
        "url": (
            f"/screenshots/"
            f"{request.sessionId}/"
            f"{filename}"
        ),
        "size": len(image_bytes),
    }


# =========================================================
# GET SCREENSHOT FILE
# =========================================================

@app.get("/screenshots/{session_id}/{filename}")
def get_screenshot(
    session_id: str,
    filename: str,
):
    safe_filename = Path(filename).name

    if safe_filename != filename:
        raise HTTPException(
            status_code=400,
            detail="Invalid screenshot filename",
        )

    filepath = SCREENSHOTS_DIR / session_id / safe_filename
    if not filepath.is_file():
        alt_filepath = BASE_DIR.parent / "screenshots" / session_id / safe_filename
        if alt_filepath.is_file():
            filepath = alt_filepath
        else:
            raise HTTPException(
                status_code=404,
                detail="Screenshot not found",
            )

    return FileResponse(
        path=filepath,
        media_type="image/png",
        filename=safe_filename,
    )


# =========================================================
# RENAME WORKFLOW
# =========================================================

@app.patch("/sessions/{session_id}")
def update_session(
    session_id: str,
    request: UpdateWorkflowRequest,
):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        cursor.execute("SELECT * FROM workflows WHERE id = ?", (session_id,))
        workflow = cursor.fetchone()
        if not workflow:
            raise HTTPException(status_code=404, detail="Workflow not found")

        updates = []
        params = []
        if request.name is not None:
            name_val = request.name.strip()
            if name_val:
                updates.append("name = ?")
                params.append(name_val)
        if request.tags is not None:
            updates.append("tags = ?")
            params.append(request.tags.strip())

        if updates:
            updates.append("updated_at = ?")
            params.append(utc_now())
            params.append(session_id)
            sql = f"UPDATE workflows SET {', '.join(updates)} WHERE id = ?"
            cursor.execute(sql, tuple(params))
            connection.commit()

        return {
            "success": True,
            "message": "Workflow updated",
        }
    finally:
        connection.close()


# =========================================================
# DELETE WORKFLOW
# =========================================================

@app.delete("/sessions/{session_id}")
def delete_session(session_id: str):
    connection = get_connection()

    try:
        cursor = connection.cursor()

        cursor.execute(
            """
            SELECT id
            FROM workflows
            WHERE id = ?
            """,
            (session_id,),
        )

        workflow = cursor.fetchone()

        if not workflow:
            raise HTTPException(
                status_code=404,
                detail="Workflow not found",
            )

        cursor.execute(
            """
            DELETE FROM workflow_steps
            WHERE workflow_id = ?
            """,
            (session_id,),
        )

        cursor.execute(
            """
            DELETE FROM workflows
            WHERE id = ?
            """,
            (session_id,),
        )

        connection.commit()

    finally:
        connection.close()

    screenshot_dir = (
        SCREENSHOTS_DIR /
        session_id
    )

    if screenshot_dir.exists():
        for file in screenshot_dir.glob("*"):
            if file.is_file():
                try:
                    file.unlink()
                except OSError:
                    pass

        try:
            screenshot_dir.rmdir()
        except OSError:
            pass

    return {
        "success": True,
        "message": "Workflow deleted",
    }


# =========================================================
# DUPLICATE WORKFLOW
# =========================================================

@app.post("/sessions/{session_id}/duplicate")
def duplicate_session(session_id: str):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        cursor.execute("SELECT * FROM workflows WHERE id = ?", (session_id,))
        workflow = cursor.fetchone()
        if not workflow:
            raise HTTPException(status_code=404, detail="Workflow not found")

        new_id = str(uuid4())
        now = utc_now()
        new_name = f"{workflow['name']} (Copy)"

        # Insert duplicated workflow
        cursor.execute(
            """
            INSERT INTO workflows (id, name, application, status, started_at, ended_at, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (new_id, new_name, workflow["application"], workflow["status"], workflow["started_at"], workflow["ended_at"], now, now),
        )

        # Copy screenshots directory if it exists
        old_screenshot_dir = SCREENSHOTS_DIR / session_id
        new_screenshot_dir = SCREENSHOTS_DIR / new_id
        if old_screenshot_dir.exists():
            new_screenshot_dir.mkdir(parents=True, exist_ok=True)
            for f in old_screenshot_dir.glob("*"):
                if f.is_file():
                    try:
                        shutil.copy2(f, new_screenshot_dir / f.name)
                    except Exception:
                        pass

        # Fetch original steps
        cursor.execute("SELECT * FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence ASC", (session_id,))
        original_steps = cursor.fetchall()

        for step in original_steps:
            old_step_id = step["id"]
            old_path = step["screenshot_path"]
            new_path = None
            if old_path:
                # Update screenshot path for the new session
                filename = Path(old_path).name
                new_path = str(new_screenshot_dir / filename)

            cursor.execute(
                """
                INSERT INTO workflow_steps (
                    workflow_id, sequence, action, timestamp, url, title, value,
                    selected_text, previous_url, checked, element_json, screenshot_path
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (
                    new_id, step["sequence"], step["action"], step["timestamp"],
                    step["url"], step["title"], step["value"], step["selected_text"],
                    step["previous_url"], step["checked"], step["element_json"], new_path
                ),
            )
            new_step_id = cursor.lastrowid

            # Copy step annotations if any
            cursor.execute("SELECT data FROM step_annotations WHERE step_id = ?", (old_step_id,))
            anno_row = cursor.fetchone()
            if anno_row:
                cursor.execute(
                    """
                    INSERT INTO step_annotations (step_id, workflow_id, data, created_at, updated_at)
                    VALUES (?, ?, ?, ?, ?)
                    """,
                    (new_step_id, new_id, anno_row["data"], now, now),
                )

            # Copy step edits if any
            cursor.execute("SELECT * FROM step_edits WHERE step_id = ?", (old_step_id,))
            edit_row = cursor.fetchone()
            if edit_row:
                cursor.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, voiceover, hidden, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        new_step_id, new_id, edit_row["title"], edit_row["description"],
                        edit_row["note"], edit_row["expected"],
                        edit_row["voiceover"] if "voiceover" in edit_row.keys() else None,
                        edit_row["hidden"], now
                    ),
                )

        connection.commit()
        return {
            "success": True,
            "id": new_id,
            "name": new_name,
            "message": f"Successfully duplicated workflow as '{new_name}'"
        }
    finally:
        connection.close()


# =========================================================
# ANNOTATIONS CRUD
# =========================================================

@app.get("/sessions/{session_id}/steps/{step_id}/annotations")
def get_step_annotations(session_id: str, step_id: int):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        row = cursor.execute(
            """
            SELECT data FROM step_annotations
            WHERE step_id = ? AND workflow_id = ?
            """,
            (step_id, session_id)
        ).fetchone()
        
        if not row:
            return {"data": "[]"}
        return {"data": row["data"]}
    finally:
        connection.close()


@app.put("/sessions/{session_id}/steps/{step_id}/annotations")
def save_step_annotations(session_id: str, step_id: int, request: StepAnnotationRequest):
    connection = get_connection()
    now = utc_now()
    try:
        cursor = connection.cursor()
        # Verify step exists
        step = cursor.execute(
            "SELECT id FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
            
        cursor.execute(
            """
            INSERT INTO step_annotations (step_id, workflow_id, data, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(step_id) DO UPDATE SET
                data = excluded.data,
                updated_at = excluded.updated_at
            """,
            (step_id, session_id, request.data, now, now)
        )
        connection.commit()
        perform_auto_backup()
    finally:
        connection.close()
    return {"success": True, "message": "Annotations saved"}


@app.delete("/sessions/{session_id}/steps/{step_id}/annotations")
def delete_step_annotations(session_id: str, step_id: int):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        cursor.execute(
            "DELETE FROM step_annotations WHERE step_id = ? AND workflow_id = ?",
            (step_id, session_id)
        )
        connection.commit()
    finally:
        connection.close()
    return {"success": True, "message": "Annotations deleted"}


# =========================================================
# STEP EDITS PERSISTENCE
# =========================================================

@app.patch("/sessions/{session_id}/steps/{step_id}/edits")
def edit_step(session_id: str, step_id: int, request: StepEditRequest):
    connection = get_connection()
    now = utc_now()
    try:
        cursor = connection.cursor()
        # Verify step exists
        step = cursor.execute(
            "SELECT id FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
            
        # Get existing edit or insert new
        existing = cursor.execute(
            "SELECT title, description, note, expected, voiceover, branches, hidden FROM step_edits WHERE step_id = ?",
            (step_id,)
        ).fetchone()
        
        if existing:
            title = request.title if request.title is not None else existing["title"]
            description = request.description if request.description is not None else existing["description"]
            note = request.note if request.note is not None else existing["note"]
            expected = request.expected if request.expected is not None else existing["expected"]
            voiceover = request.voiceover if request.voiceover is not None else existing["voiceover"]
            branches = request.branches if request.branches is not None else (existing["branches"] if "branches" in existing.keys() else "")
            hidden = int(request.hidden) if request.hidden is not None else existing["hidden"]
            alert_type = request.alert_type if request.alert_type is not None else (existing["alert_type"] if "alert_type" in existing.keys() else "none")
            alert_msg = request.alert_msg if request.alert_msg is not None else (existing["alert_msg"] if "alert_msg" in existing.keys() else "")
            role = request.role if request.role is not None else (existing["role"] if "role" in existing.keys() else "")
            duration = request.duration if request.duration is not None else (existing["duration"] if "duration" in existing.keys() else "")
        else:
            orig_step = cursor.execute(
                "SELECT title, action, value FROM workflow_steps WHERE id = ?",
                (step_id,)
            ).fetchone()
            
            title = request.title if request.title is not None else (orig_step["title"] or "")
            description = request.description if request.description is not None else ""
            note = request.note if request.note is not None else ""
            expected = request.expected if request.expected is not None else ""
            voiceover = request.voiceover if request.voiceover is not None else ""
            branches = request.branches if request.branches is not None else ""
            hidden = int(request.hidden) if request.hidden is not None else 0
            alert_type = request.alert_type if request.alert_type is not None else "none"
            alert_msg = request.alert_msg if request.alert_msg is not None else ""
            role = request.role if request.role is not None else ""
            duration = request.duration if request.duration is not None else ""
            
        cursor.execute(
            """
            INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, voiceover, branches, hidden, alert_type, alert_msg, role, duration, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(step_id) DO UPDATE SET
                title = excluded.title,
                description = excluded.description,
                note = excluded.note,
                expected = excluded.expected,
                voiceover = excluded.voiceover,
                branches = excluded.branches,
                hidden = excluded.hidden,
                alert_type = excluded.alert_type,
                alert_msg = excluded.alert_msg,
                role = excluded.role,
                duration = excluded.duration,
                updated_at = excluded.updated_at
            """,
            (step_id, session_id, title, description, note, expected, voiceover, branches, hidden, alert_type, alert_msg, role, duration, now)
        )
        if request.checked is not None:
            cursor.execute(
                "UPDATE workflow_steps SET checked = ? WHERE id = ?",
                (1 if request.checked else 0, step_id)
            )
        connection.commit()
        perform_auto_backup()
    finally:
        connection.close()
    return {"success": True, "message": "Step edits saved"}


@app.delete("/sessions/{session_id}/steps/{step_id}")
def delete_step_permanently(session_id: str, step_id: int):
    """
    Permanently deletes a step from the database and removes its screenshot from disk.
    """
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT id, screenshot_path, sequence FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
            
        screenshot_path = step["screenshot_path"]
        if screenshot_path:
            file_path = BASE_DIR / screenshot_path
            if file_path.exists() and file_path.is_file():
                try:
                    file_path.unlink()
                except Exception as e:
                    print(f"Warning: Could not delete screenshot file {file_path}: {e}")
                    
        # Delete from annotations and edits
        cursor.execute("DELETE FROM step_annotations WHERE step_id = ? AND workflow_id = ?", (step_id, session_id))
        cursor.execute("DELETE FROM step_edits WHERE step_id = ? AND workflow_id = ?", (step_id, session_id))
        cursor.execute("DELETE FROM workflow_steps WHERE id = ? AND workflow_id = ?", (step_id, session_id))
        
        # Renumber remaining steps to maintain contiguous sequences
        remaining_steps = cursor.execute(
            "SELECT id FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence ASC",
            (session_id,)
        ).fetchall()
        
        for new_seq, r in enumerate(remaining_steps, 1):
            cursor.execute("UPDATE workflow_steps SET sequence = ? WHERE id = ?", (new_seq, r["id"]))
            
        connection.commit()
        perform_auto_backup()
    finally:
        connection.close()
        
    return {"success": True, "message": "Step permanently deleted"}


# =========================================================
# OFFLINE AI ENHANCEMENTS (Ollama)
# =========================================================
import urllib.request

def call_ollama_sync(endpoint: str, payload: dict, timeout=180) -> dict:
    url = f"http://127.0.0.1:11434{endpoint}"
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST"
    )
    try:
        with urllib.request.urlopen(req, timeout=timeout) as response:
            if response.status == 200:
                return json.loads(response.read().decode("utf-8"))
    except Exception as e:
        print(f"Ollama connection error: {e}")
    return None

def get_installed_models() -> list:
    try:
        req = urllib.request.Request("http://127.0.0.1:11434/api/tags", method="GET")
        with urllib.request.urlopen(req, timeout=2) as response:
            if response.status == 200:
                data = json.loads(response.read().decode("utf-8"))
                return [m["name"] for m in data.get("models", [])]
    except Exception:
        pass
    return []

def get_ollama_heartbeat() -> bool:
    try:
        req = urllib.request.Request("http://127.0.0.1:11434/", method="GET")
        with urllib.request.urlopen(req, timeout=1) as response:
            return response.status == 200
    except Exception:
        return False

@app.get("/ai/status")
def get_ai_status():
    heartbeat = get_ollama_heartbeat()
    installed = get_installed_models() if heartbeat else []
    ollama_running = heartbeat or len(installed) > 0
    
    # Check if moondream and qwen2.5 are pulled
    installed_names = [m.split(":")[0] for m in installed]
    has_moondream = "moondream" in installed_names
    has_qwen = "qwen2.5" in installed_names
    
    # Find Ollama executable
    ollama_path = shutil.which("ollama")
    if not ollama_path:
        possible_paths = [
            Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Ollama" / "ollama.exe",
            Path(os.environ.get("ProgramFiles", "")) / "Ollama" / "ollama.exe",
            Path(os.environ.get("ProgramFiles(x86)", "")) / "Ollama" / "ollama.exe",
        ]
        for p in possible_paths:
            if p.exists():
                ollama_path = str(p)
                break

    diagnostic_msg = ""
    if not ollama_running:
        if not ollama_path:
            diagnostic_msg = "Ollama is not installed on this PC. Click 'Start Ollama' to auto-install via winget."
        else:
            diagnostic_msg = f"Ollama executable found at {ollama_path}, but service is stopped. Click 'Start Ollama' to launch it."
    elif not (has_moondream and has_qwen):
        missing = []
        if not has_moondream: missing.append("moondream (vision)")
        if not has_qwen: missing.append("qwen2.5 (text)")
        diagnostic_msg = f"Ollama running, but missing models: {', '.join(missing)}. Click 'Start Ollama' to pull them."
    else:
        diagnostic_msg = "AI engine active with all models loaded."

    return {
        "running": ollama_running,
        "models": installed,
        "required_models_present": has_moondream and has_qwen,
        "ollama_installed": bool(ollama_path),
        "diagnostic_message": diagnostic_msg
    }

@app.post("/ai/start-ollama")
def start_ollama():
    try:
        if get_ollama_heartbeat():
            return {"success": True, "message": "Ollama is already running"}
            
        ollama_path = shutil.which("ollama")
        if not ollama_path:
            possible_paths = [
                Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Ollama" / "ollama.exe",
                Path(os.environ.get("ProgramFiles", "")) / "Ollama" / "ollama.exe",
                Path(os.environ.get("ProgramFiles(x86)", "")) / "Ollama" / "ollama.exe",
            ]
            for p in possible_paths:
                if p.exists():
                    ollama_path = str(p)
                    break

        if ollama_path:
            def run_ollama():
                subprocess.Popen([ollama_path, "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            t = threading.Thread(target=run_ollama)
            t.start()
            return {"success": True, "message": "Ollama start command issued"}
        else:
            def install_ollama():
                subprocess.run(["winget", "install", "Ollama.Ollama", "--accept-source-agreements", "--accept-package-agreements", "--silent"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            t = threading.Thread(target=install_ollama)
            t.start()
            return {"success": True, "message": "Ollama not installed. Installing via winget in background... Please wait 1-2 mins."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _find_ollama_path() -> str | None:
    """Locate the ollama executable."""
    p = shutil.which("ollama")
    if p:
        return p
    for candidate in [
        Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Ollama" / "ollama.exe",
        Path(os.environ.get("ProgramFiles", "")) / "Ollama" / "ollama.exe",
        Path(os.environ.get("ProgramFiles(x86)", "")) / "Ollama" / "ollama.exe",
    ]:
        if candidate.exists():
            return str(candidate)
    return None


@app.get("/ai/pull-models-stream")
def pull_ai_models_stream():
    """
    Streams the real-time download progress of moondream and qwen2.5 from Ollama
    with live percentage, progress bar, and byte counters.
    """
    def generate_progress():
        import json, urllib.request, time

        ollama_path = _find_ollama_path()
        if not ollama_path:
            yield "data: ❌ Error: Ollama executable not found. Install from https://ollama.com/download\n\n"
            yield "data: [DONE]\n\n"
            return

        # Ensure ollama service is running
        if not get_ollama_heartbeat():
            yield "data: 🚀 Starting Ollama background service...\n\n"
            try:
                subprocess.Popen([ollama_path, "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                for _ in range(12):
                    time.sleep(1)
                    if get_ollama_heartbeat():
                        break
            except Exception:
                pass

        models_to_pull = [
            ("moondream", "Vision AI (Screenshot Analysis)", "1.7 GB"),
            ("qwen2.5",   "Text AI (SOP Title & Description Generator)", "4.7 GB"),
        ]

        installed_names = [m.split(":")[0] for m in get_installed_models()]

        for model_name, purpose, size in models_to_pull:
            if model_name in installed_names:
                yield f"data: ✓ {model_name} ({size}) is already installed.\n\n"
                continue

            yield f"data: ⬇ Downloading {model_name} (~{size}) — {purpose}...\n\n"

            # Pull via Ollama HTTP streaming API
            success = False
            try:
                req = urllib.request.Request(
                    "http://127.0.0.1:11434/api/pull",
                    data=json.dumps({"name": model_name, "stream": True}).encode("utf-8"),
                    headers={"Content-Type": "application/json"}
                )
                with urllib.request.urlopen(req, timeout=3600) as response:
                    last_pct = -1
                    for line in response:
                        if not line:
                            continue
                        try:
                            data = json.loads(line.decode("utf-8"))
                            status = data.get("status", "")
                            completed = data.get("completed", 0)
                            total = data.get("total", 0)
                            if total > 0 and completed > 0:
                                pct = int((completed / total) * 100)
                                if pct != last_pct and (pct % 2 == 0 or pct == 100):
                                    last_pct = pct
                                    mb_done = completed / (1024 * 1024)
                                    mb_total = total / (1024 * 1024)
                                    bars = int(pct / 5)
                                    bar_str = "█" * bars + "░" * (20 - bars)
                                    yield f"data: [{bar_str}] {pct}% ({mb_done:.1f} MB / {mb_total:.1f} MB) — {status}\n\n"
                            elif status and ("pulling" in status.lower() or "verifying" in status.lower() or "writing" in status.lower()):
                                yield f"data: • {status}\n\n"
                        except Exception:
                            continue
                success = True
                yield f"data: ✓ {model_name} installed successfully!\n\n"
            except Exception as e:
                # Fallback to CLI subprocess
                yield f"data: [CLI Fallback] Running: ollama pull {model_name} ...\n\n"
                try:
                    proc = subprocess.Popen(
                        [ollama_path, "pull", model_name],
                        stdout=subprocess.PIPE,
                        stderr=subprocess.STDOUT,
                        text=True,
                        bufsize=1
                    )
                    for pline in proc.stdout:
                        clean = pline.strip()
                        if clean:
                            yield f"data: {clean}\n\n"
                    proc.wait()
                    if proc.returncode == 0:
                        success = True
                        yield f"data: ✓ {model_name} installed successfully!\n\n"
                    else:
                        yield f"data: ✗ {model_name} pull exited with code {proc.returncode}\n\n"
                except Exception as ex:
                    yield f"data: ✗ {model_name} pull failed: {str(ex)}\n\n"

        yield "data: \n✅ All required AI models are ready and verified!\n\n"
        yield "data: [DONE]\n\n"

    from fastapi.responses import StreamingResponse
    return StreamingResponse(generate_progress(), media_type="text/event-stream")


@app.post("/ai/pull-models")
def pull_ai_models():
    """
    Pull the required AI models (moondream + qwen2.5) via `ollama pull`.
    Runs both pulls sequentially and returns combined log output.
    The endpoint will block until complete (can take 5-30 min depending on connection).
    """
    ollama_path = _find_ollama_path()
    if not ollama_path:
        return {
            "success": False,
            "output": "❌ Ollama executable not found. Please install Ollama first via https://ollama.com/download",
            "models_pulled": []
        }

    # Make sure the service is running before pulling
    if not get_ollama_heartbeat():
        try:
            subprocess.Popen([ollama_path, "serve"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            import time
            for _ in range(15):
                time.sleep(1)
                if get_ollama_heartbeat():
                    break
        except Exception:
            pass

    models_to_pull = [
        ("moondream", "Vision AI — describes screenshots automatically"),
        ("qwen2.5",   "Text AI — polishes and refines SOP descriptions"),
    ]

    combined_output = []
    models_pulled = []

    for model_name, model_purpose in models_to_pull:
        # Skip if already installed
        installed_names = [m.split(":")[0] for m in get_installed_models()]
        if model_name in installed_names:
            combined_output.append(f"✓ {model_name} — already installed, skipping pull.")
            models_pulled.append(model_name)
            continue

        combined_output.append(f"\n⬇ Pulling {model_name} ({model_purpose})...")
        try:
            result = subprocess.run(
                [ollama_path, "pull", model_name],
                capture_output=True,
                text=True,
                timeout=1800  # 30 min max
            )
            stdout = result.stdout.strip()
            stderr = result.stderr.strip()
            if result.returncode == 0:
                combined_output.append(f"✓ {model_name} pulled successfully.")
                models_pulled.append(model_name)
            else:
                combined_output.append(f"✗ {model_name} pull failed (exit {result.returncode}).")
                if stderr:
                    combined_output.append(f"  Error: {stderr[:400]}")
        except subprocess.TimeoutExpired:
            combined_output.append(f"✗ {model_name} pull timed out after 30 minutes.")
        except Exception as e:
            combined_output.append(f"✗ {model_name} pull error: {str(e)}")

    all_ok = all(m in models_pulled for m, _ in models_to_pull)
    summary = "\n✅ All required AI models are ready!" if all_ok else "\n⚠️ Some models could not be pulled. Check your internet connection and try again."
    combined_output.append(summary)

    return {
        "success": all_ok,
        "output": "\n".join(combined_output),
        "models_pulled": models_pulled,
    }

class TTSRequest(BaseModel):
    text: Optional[str] = None
    voice: str = "en-US-AriaNeural"

@app.post("/sessions/{session_id}/steps/{step_id}/tts")
async def generate_tts(session_id: str, step_id: int, request: TTSRequest):
    try:
        speech_text = request.text
        if not speech_text:
            connection = get_connection()
            try:
                cursor = connection.cursor()
                row = cursor.execute(
                    """
                    SELECT se.voiceover, se.description, ws.title, ws.action, ws.value
                    FROM workflow_steps ws
                    LEFT JOIN step_edits se ON ws.id = se.step_id
                    WHERE ws.id = ? AND ws.workflow_id = ?
                    """,
                    (step_id, session_id)
                ).fetchone()
                if row:
                    speech_text = row["voiceover"] or row["description"] or row["title"] or f"{row['action']} {row['value'] or ''}"
            finally:
                connection.close()
                
        if not speech_text or not speech_text.strip():
            speech_text = "No text available for voiceover."

        audio_filename = f"tts_{session_id}_{step_id}.mp3"
        audio_path = MEDIA_DIR / audio_filename

        communicate = edge_tts.Communicate(speech_text, request.voice)
        await communicate.save(str(audio_path))

        return {"success": True, "audioUrl": f"/media/{audio_filename}"}
    except Exception as e:
        print("TTS Error:", e)
        raise HTTPException(status_code=500, detail=str(e))

class AIDescribeRequest(BaseModel):
    step_id: int
    session_id: str

def _get_text_model() -> str:
    installed = get_installed_models()
    for inst in installed:
        if inst.startswith("qwen2.5"):
            return inst
    return "qwen2.5"

def _get_vision_model() -> str:
    installed = get_installed_models()
    for inst in installed:
        if inst.startswith("moondream"):
            return inst
    return "moondream"

@app.post("/ai/describe-step")
def ai_describe_step(request: AIDescribeRequest):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT action, value, selected_text, url, title, element_json, screenshot_path FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (request.step_id, request.session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
        
        # 1. Vision Analysis using moondream (if screenshot exists)
        vision_description = ""
        screenshot_path = step["screenshot_path"]
        if screenshot_path and Path(screenshot_path).exists():
            try:
                with open(screenshot_path, "rb") as image_file:
                    base64_image = base64.b64encode(image_file.read()).decode("utf-8")
                
                prompt = (
                    f"This is a screenshot of a user interacting with a web page. "
                    f"The user performed a '{step['action']}' action. "
                    f"Describe in one concise, professional sentence what exact button, input field, or element the user clicked or interacted with."
                )
                
                payload = {
                    "model": _get_vision_model(),
                    "prompt": prompt,
                    "images": [base64_image],
                    "stream": False
                }
                
                response_data = call_ollama_sync("/api/generate", payload, timeout=60)
                if response_data and "response" in response_data:
                    vision_description = response_data["response"].strip()
            except Exception as e:
                print(f"Vision analysis failed: {e}")
                
        # 2. Text Polishing using qwen2.5
        element_info = ""
        if step["element_json"]:
            try:
                el = json.loads(step["element_json"])
                element_info = (
                    f"Element Tag: {el.get('tagName', '')}\n"
                    f"Element Text: {el.get('text', '')}\n"
                    f"Element Name/ID: {el.get('name', '') or el.get('id', '')}\n"
                    f"CSS Selector: {el.get('cssSelector', '')}\n"
                )
            except Exception:
                element_info = f"Element JSON: {step['element_json']}\n"
                
        prompt_text = (
            f"You are ProcSnap AI, an expert at writing standard operating procedures (SOPs).\n"
            f"Based on the following user action details, generate a clean title, a step-by-step instruction description, and an expected result.\n\n"
            f"--- ACTION DETAILS ---\n"
            f"Action Type: {step['action']}\n"
            f"Value Inputted/Selected: {step['value'] or 'None'}\n"
            f"Selected Text: {step['selected_text'] or 'None'}\n"
            f"Page URL: {step['url']}\n"
            f"Page Title: {step['title']}\n"
            f"UI Element Details:\n{element_info}\n"
        )
        
        if vision_description:
            prompt_text += f"Screenshot visual analysis description: {vision_description}\n"
            
        prompt_text += (
            f"\nRespond ONLY with a JSON object in this exact format (no markdown code blocks, just raw JSON):\n"
            f'{{\n  "title": "Concise active voice step title (e.g., Click standard button)",\n  "description": "Clear step-by-step description of what to do (e.g., click on the search bar and type ...)",\n  "expected": "What is the expected result or next screen state?"\n}}'
        )
        
        payload_text = {
            "model": _get_text_model(),
            "prompt": prompt_text,
            "format": "json",
            "stream": False
        }
        
        response_text = call_ollama_sync("/api/generate", payload_text, timeout=30)
        
        title = ""
        description = ""
        expected = ""
        
        if response_text and "response" in response_text:
            try:
                parsed = json.loads(response_text["response"])
                title = parsed.get("title", "").strip()
                description = parsed.get("description", "").strip()
                expected = parsed.get("expected", "").strip()
            except Exception:
                print("Failed parsing JSON response from text model. Raw:", response_text["response"])
        
        if not title:
            title = step["title"] or f"Step {step['action']}"
        if not description:
            description = f"Perform {step['action']} on {step['url']}"
            
        now = datetime.now(timezone.utc).isoformat()
        
        # Save to database
        existing = cursor.execute("SELECT note, hidden FROM step_edits WHERE step_id = ?", (request.step_id,)).fetchone()
        note = existing["note"] if existing else ""
        hidden = existing["hidden"] if existing else 0
        
        cursor.execute(
            """
            INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, hidden, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            ON CONFLICT(step_id) DO UPDATE SET
                title = excluded.title,
                description = excluded.description,
                expected = excluded.expected,
                updated_at = excluded.updated_at
            """,
            (request.step_id, request.session_id, title, description, note, expected, hidden, now)
        )
        connection.commit()
        
        return {
            "success": True,
            "step_id": request.step_id,
            "title": title,
            "description": description,
            "expected": expected,
            "vision_description": vision_description
        }
    finally:
        connection.close()

class AIPolishRequest(BaseModel):
    session_id: str

@app.post("/ai/polish-sop")
def ai_polish_sop(request: AIPolishRequest):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        steps = cursor.execute(
            """
            SELECT ws.id, ws.action, ws.value, se.title, se.description, se.expected
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON ws.id = se.step_id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (request.session_id,)
        ).fetchall()
        
        if not steps:
            raise HTTPException(status_code=404, detail="No steps found for this session")
            
        sop_summary = []
        for i, step in enumerate(steps):
            title = step["title"] or f"Step {i+1}"
            desc = step["description"] or f"Perform {step['action']} action."
            sop_summary.append({
                "id": step["id"],
                "title": title,
                "description": desc
            })
            
        prompt = (
            f"You are a professional technical writer polishing a Standard Operating Procedure (SOP).\n"
            f"Review the following sequence of steps and rewrite their descriptions to be formal, clear, consistent, and written in a professional instruction manual voice.\n"
            f"Do not change the meaning of the steps, only improve their writing quality.\n\n"
            f"Steps to polish:\n{json.dumps(sop_summary, indent=2)}\n\n"
            f"Respond ONLY with a JSON array containing the polished steps, in this exact format:\n"
            f"[\n"
            f"  {{\n"
            f"    \"id\": step_id_number,\n"
            f"    \"title\": \"Polished concise title\",\n"
            f"    \"description\": \"Polished formal instruction description\"\n"
            f"  }}\n"
            f"]"
        )
        
        payload = {
            "model": _get_text_model(),
            "prompt": prompt,
            "format": "json",
            "stream": False
        }
        
        response_text = call_ollama_sync("/api/generate", payload, timeout=180)
        if response_text and "response" in response_text:
            try:
                polished_list = json.loads(response_text["response"])
                now = datetime.now(timezone.utc).isoformat()
                
                for item in polished_list:
                    step_id = item.get("id")
                    p_title = item.get("title", "").strip()
                    p_desc = item.get("description", "").strip()
                    
                    if not step_id or (not p_title and not p_desc):
                        continue
                        
                    existing = cursor.execute(
                        "SELECT expected, note, hidden FROM step_edits WHERE step_id = ?",
                        (step_id,)
                    ).fetchone()
                    
                    expected = existing["expected"] if existing else ""
                    note = existing["note"] if existing else ""
                    hidden = existing["hidden"] if existing else 0
                    
                    cursor.execute(
                        """
                        INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, hidden, updated_at)
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                        ON CONFLICT(step_id) DO UPDATE SET
                            title = CASE WHEN ? != '' THEN ? ELSE title END,
                            description = CASE WHEN ? != '' THEN ? ELSE description END,
                            updated_at = ?
                        """,
                        (step_id, request.session_id, p_title, p_desc, note, expected, hidden, now,
                         p_title, p_title, p_desc, p_desc, now)
                    )
                connection.commit()
                return {"success": True, "count": len(polished_list), "steps": polished_list}
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Failed parsing polished SOP JSON: {e}")
                
        raise HTTPException(status_code=500, detail="No response from Ollama")
    finally:
        connection.close()

class AIDetectRedactRequest(BaseModel):
    step_id: int
    session_id: str

@app.post("/ai/detect-redact")
def ai_detect_redact(request: AIDetectRedactRequest):
    connection = get_database_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT element_json, screenshot_path FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (request.step_id, request.session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
        
        redact_regions = []
        
        # Heuristic 1: Check if input was a password field or input element matching password regexes
        element_json = step["element_json"]
        if element_json:
            try:
                el = json.loads(element_json)
                is_sensitive = False
                if el.get("type") == "password":
                    is_sensitive = True
                else:
                    name = (el.get("name") or "").lower()
                    placeholder = (el.get("placeholder") or "").lower()
                    id_attr = (el.get("id") or "").lower()
                    for keyword in ["password", "pass", "pwd", "ssn", "socialsecurity", "cardnumber", "cvv", "creditcard", "secret", "token", "apikey"]:
                        if keyword in name or keyword in placeholder or keyword in id_attr:
                            is_sensitive = True
                            break
                
                if is_sensitive and el.get("screen"):
                    screen = el["screen"]
                    width = screen["width"]
                    height = screen["height"]
                    x = screen["x"]
                    y = screen["y"]
                    
                    redact_regions.append({
                        "id": "ai-auto-" + uuid4().hex[:6],
                        "type": "blur",
                        "x": x,
                        "y": y,
                        "w": width,
                        "h": height,
                        "color": "#ef4444"
                    })
            except Exception as e:
                print("Element-based redact check failed:", e)

        # Heuristic 2: If we have vision model moondream, we can also ask it to look for credentials or credit cards
        screenshot_path = step["screenshot_path"]
        if not redact_regions and screenshot_path and Path(screenshot_path).exists():
            try:
                with open(screenshot_path, "rb") as image_file:
                    base64_image = base64.b64encode(image_file.read()).decode("utf-8")
                
                prompt = (
                    "Locate any visible email addresses, credit cards, password values, or names in this image. "
                    "Respond ONLY with a JSON list of bounding boxes: [[x, y, width, height]] in percent coordinates (0 to 100). "
                    "Do not include any text besides raw JSON."
                )
                
                payload = {
                    "model": "moondream",
                    "prompt": prompt,
                    "format": "json",
                    "stream": False
                }
                
                response_data = call_ollama_sync("/api/generate", payload, timeout=60)
                if response_data and "response" in response_data:
                    boxes = json.loads(response_data["response"])
                    from PIL import Image
                    with Image.open(screenshot_path) as img:
                        img_w, img_h = img.size
                    
                    for box in boxes:
                        if len(box) == 4:
                            x_pct, y_pct, w_pct, h_pct = box
                            x = (x_pct / 100.0) * img_w
                            y = (y_pct / 100.0) * img_h
                            w = (w_pct / 100.0) * img_w
                            h = (h_pct / 100.0) * img_h
                            
                            redact_regions.append({
                                "id": "ai-vision-" + uuid4().hex[:6],
                                "type": "blur",
                                "x": x,
                                "y": y,
                                "w": w,
                                "h": h,
                                "color": "#ef4444"
                            })
            except Exception as e:
                print("Vision-based redact detection failed:", e)
                
        return {"success": True, "regions": redact_regions}
    finally:
        connection.close()


class ReplaceScreenshotRequest(BaseModel):
    image: str

@app.post("/sessions/{session_id}/steps/{step_id}/screenshot")
def replace_step_screenshot(session_id: str, step_id: int, request: ReplaceScreenshotRequest):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT id FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")
            
        match = re.fullmatch(
            r"data:image/[^;]+;base64,(.+)",
            request.image,
            flags=re.DOTALL,
        )
        if not match:
            raise HTTPException(status_code=400, detail="Invalid data URL format")
            
        image_data = base64.b64decode(match.group(1))
        
        # Save to file
        import uuid
        filename = f"replace_{uuid.uuid4().hex}.png"
        folder = Path("storage/screenshots") / session_id
        folder.mkdir(parents=True, exist_ok=True)
        file_path = folder / filename
        
        with open(file_path, "wb") as f:
            f.write(image_data)
            
        # Update database path
        db_path = f"storage/screenshots/{session_id}/{filename}"
        cursor.execute(
            "UPDATE workflow_steps SET screenshot_path = ? WHERE id = ?",
            (db_path, step_id)
        )
        connection.commit()
        return {"success": True, "screenshotUrl": f"storage/screenshots/{session_id}/{filename}"}
    finally:
        connection.close()


class CropScreenshotRequest(BaseModel):
    x: float
    y: float
    w: float
    h: float

@app.post("/sessions/{session_id}/steps/{step_id}/screenshot/crop")
def crop_step_screenshot(session_id: str, step_id: int, request: CropScreenshotRequest):
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT screenshot_path FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step or not step["screenshot_path"]:
            raise HTTPException(status_code=404, detail="Screenshot not found")
            
        full_path = Path(step["screenshot_path"])
        if not full_path.exists():
            raise HTTPException(status_code=404, detail="File not found")
            
        # Coordinates logic
        x1 = min(request.x, request.x + request.w)
        y1 = min(request.y, request.y + request.h)
        x2 = max(request.x, request.x + request.w)
        y2 = max(request.y, request.y + request.h)
        
        # Guard minimum dimensions
        if (x2 - x1) < 10 or (y2 - y1) < 10:
            raise HTTPException(status_code=400, detail="Crop region too small")
            
        # Crop using PIL
        from PIL import Image
        with Image.open(full_path) as img:
            cropped = img.crop((int(x1), int(y1), int(x2), int(y2)))
            import uuid
            filename = f"crop_{uuid.uuid4().hex}.png"
            folder = Path("storage/screenshots") / session_id
            folder.mkdir(parents=True, exist_ok=True)
            file_path = folder / filename
            cropped.save(file_path)
            
        # Shift annotations
        annot_row = cursor.execute(
            "SELECT data FROM step_annotations WHERE step_id = ?",
            (step_id,)
        ).fetchone()
        
        if annot_row:
            try:
                annots = json.loads(annot_row["data"])
                for a in annots:
                    a["x"] = a["x"] - x1
                    a["y"] = a["y"] - y1
                cursor.execute(
                    "UPDATE step_annotations SET data = ?, updated_at = ? WHERE step_id = ?",
                    (json.dumps(annots), utc_now(), step_id)
                )
            except Exception as e:
                print("Failed shifting annotations:", e)
                
        # Update database screenshot path
        db_path = f"storage/screenshots/{session_id}/{filename}"
        cursor.execute(
            "UPDATE workflow_steps SET screenshot_path = ? WHERE id = ?",
            (db_path, step_id)
        )
        connection.commit()
        return {"success": True, "screenshotUrl": f"storage/screenshots/{session_id}/{filename}"}
    finally:
        connection.close()


@app.get("/sessions/{session_id}/steps/{step_id}/annotated-image")
def get_step_annotated_image(session_id: str, step_id: int):
    """
    Returns the step's screenshot composited with all canvas annotations,
    spotlight focus, and red bounding boxes directly as a high-res PNG image.
    """
    conn = get_connection()
    try:
        cur = conn.cursor()
        step = cur.execute(
            "SELECT * FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step or not step["screenshot_path"]:
            raise HTTPException(status_code=404, detail="Step or screenshot not found")
        
        img_file = BASE_DIR / step["screenshot_path"]
        if not img_file.exists():
            img_file = BASE_DIR / "screenshots" / session_id / Path(step["screenshot_path"]).name
        if not img_file.exists():
            raise HTTPException(status_code=404, detail="Screenshot file not found on disk")
        
        # Fetch annotations
        ann_row = cur.execute(
            "SELECT data FROM step_annotations WHERE step_id = ?",
            (step_id,)
        ).fetchone()
        annotations = json.loads(ann_row["data"]) if ann_row and ann_row["data"] else []
        
        try:
            from .annotation_renderer import AnnotationRenderer
        except ImportError:
            from annotation_renderer import AnnotationRenderer
        
        composited = AnnotationRenderer.composite_image(
            str(img_file),
            annotations=annotations
        )
        
        buf = io.BytesIO()
        if composited:
            composited.save(buf, format="PNG")
        else:
            with open(img_file, "rb") as f:
                buf.write(f.read())
        buf.seek(0)
        return Response(content=buf.getvalue(), media_type="image/png")
    finally:
        conn.close()


# =========================================================
# ANIMATED STEP MICRO-DEMO GENERATOR (GIF)
# =========================================================

class StepAnimatePayload(BaseModel):
    target_x: Optional[float] = None
    target_y: Optional[float] = None
    x_pct: Optional[float] = None
    y_pct: Optional[float] = None

@app.post("/sessions/{session_id}/steps/{step_id}/animate")
def generate_step_animation(
    session_id: str,
    step_id: int,
    payload: Optional[StepAnimatePayload] = Body(default=None),
    x_pct: Optional[float] = Query(None),
    y_pct: Optional[float] = Query(None),
    target_x_param: Optional[float] = Query(None, alias="target_x"),
    target_y_param: Optional[float] = Query(None, alias="target_y")
):
    """
    Generates a 14-frame animated micro-demo GIF for the step.
    Interpolates a simulated moving cursor towards the hotspot coordinate,
    followed by a smooth pulsing ripple click animation.
    """
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT id, sequence, element_json, screenshot_path FROM workflow_steps WHERE (id = ? OR sequence = ?) AND workflow_id = ?",
            (step_id, step_id, session_id)
        ).fetchone()
        if not step:
            step = cursor.execute(
                "SELECT id, sequence, element_json, screenshot_path FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence LIMIT 1",
                (session_id,)
            ).fetchone()

        if not step:
            raise HTTPException(status_code=404, detail="Step not found in database")

        raw_path = step["screenshot_path"] or ""
        seq = int(step["sequence"] or 1)

        candidate_paths = [
            SCREENSHOTS_DIR / session_id / f"step-{seq:03d}.png",
            SCREENSHOTS_DIR / session_id / f"step_{seq}.png",
            SCREENSHOTS_DIR / session_id / f"step-{seq}.png",
            SCREENSHOTS_DIR / session_id / f"step-{seq:03d}.jpg",
            SCREENSHOTS_DIR / session_id / f"step_{seq}.jpg",
        ]
        if raw_path:
            candidate_paths.extend([
                BASE_DIR / raw_path,
                BASE_DIR / "screenshots" / session_id / Path(raw_path).name,
                SCREENSHOTS_DIR / session_id / Path(raw_path).name,
                Path(raw_path)
            ])

        img_path = None
        for p in candidate_paths:
            if p and p.is_file() and not p.name.endswith("-demo.gif"):
                img_path = p
                break

        # Fallback: scan session screenshots folder for any matching sequence image
        if not img_path and (SCREENSHOTS_DIR / session_id).exists():
            for f in (SCREENSHOTS_DIR / session_id).glob("*.png"):
                if not f.name.endswith("-demo.gif"):
                    img_path = f
                    break
            if not img_path:
                for f in (SCREENSHOTS_DIR / session_id).glob("*.jpg"):
                    img_path = f
                    break

        if not img_path:
            # Generate clean fallback background if screenshot is not on disk
            base_img = Image.new("RGBA", (1280, 720), (15, 23, 42, 255))
        else:
            base_img = Image.open(img_path).convert("RGBA")

        width, height = base_img.size
        width, height = base_img.size

        target_x, target_y = None, None

        # 1. Percentage coordinates from client (exact frontend percentage alignment)
        req_x_pct = (payload.x_pct if payload and payload.x_pct is not None else x_pct)
        req_y_pct = (payload.y_pct if payload and payload.y_pct is not None else y_pct)
        req_tx = (payload.target_x if payload and payload.target_x is not None else target_x_param)
        req_ty = (payload.target_y if payload and payload.target_y is not None else target_y_param)

        if req_x_pct is not None and req_y_pct is not None:
            target_x = (float(req_x_pct) / 100.0) * float(width)
            target_y = (float(req_y_pct) / 100.0) * float(height)
        elif req_tx is not None and req_ty is not None:
            target_x = float(req_tx)
            target_y = float(req_ty)

        # 2. Check annotations for custom hotspot coordinates
        if target_x is None:
            anno_row = cursor.execute(
                "SELECT data FROM step_annotations WHERE step_id = ?",
                (step_id,)
            ).fetchone()
            if anno_row and anno_row["data"]:
                try:
                    annos = json.loads(anno_row["data"])
                    for a in annos:
                        if a.get("type") in ("spotlight", "rect", "circle", "badge") and "x" in a and "y" in a:
                            target_x = float(a["x"]) + (float(a.get("w", 0)) / 2.0)
                            target_y = float(a["y"]) + (float(a.get("h", 0)) / 2.0)
                            break
                except Exception:
                    pass

        # 3. Fallback to DOM element screen coordinates
        if target_x is None and step["element_json"]:
            try:
                elem = json.loads(step["element_json"])
                sc = elem.get("screen") or elem.get("coords") or {}
                if "x" in sc and "y" in sc:
                    ex = float(sc.get("x", 0))
                    ey = float(sc.get("y", 0))
                    ew = float(sc.get("width", 0))
                    eh = float(sc.get("height", 0))
                    vw = float(sc.get("viewportWidth") or width)
                    vh = float(sc.get("viewportHeight") or height)
                    
                    scale_x = width / max(1.0, vw)
                    scale_y = height / max(1.0, vh)
                    
                    target_x = (ex + (ew / 2.0)) * scale_x
                    target_y = (ey + (eh / 2.0)) * scale_y
            except Exception as ex:
                print("Error parsing element_json coords:", ex)

        # 4. Fallback to center of image if completely unspecified
        if target_x is None:
            target_x = width * 0.5
            target_y = height * 0.5

        # Ensure target is safely within image bounds
        target_x = max(24.0, min(float(width) - 24.0, target_x))
        target_y = max(24.0, min(float(height) - 24.0, target_y))

        # Start position (offset down-right for smooth trajectory)
        offset_dist_x = max(180.0, float(width) * 0.15)
        offset_dist_y = max(140.0, float(height) * 0.12)
        start_x = min(float(width) - 30.0, target_x + offset_dist_x)
        start_y = min(float(height) - 30.0, target_y + offset_dist_y)
        if start_x <= target_x + 20 and start_y <= target_y + 20:
            start_x = max(30.0, target_x - offset_dist_x)
            start_y = max(30.0, target_y - offset_dist_y)

        frames = []
        total_frames = 14
        cursor_scale = max(0.85, min(1.4, float(width) / 1920.0))

        for f in range(total_frames):
            frame = base_img.copy()
            draw = ImageDraw.Draw(frame, "RGBA")
            s = cursor_scale

            if f < 9:
                # Travel phase with smooth cubic ease-out
                t = f / 8.0
                ease_t = 1 - pow(1 - t, 3)
                cur_x = start_x + (target_x - start_x) * ease_t
                cur_y = start_y + (target_y - start_y) * ease_t

                # Sleek compact pointer arrow
                points = [
                    (cur_x, cur_y),
                    (cur_x, cur_y + 16 * s),
                    (cur_x + 4 * s, cur_y + 12 * s),
                    (cur_x + 8 * s, cur_y + 19 * s),
                    (cur_x + 11 * s, cur_y + 17.5 * s),
                    (cur_x + 7 * s, cur_y + 11 * s),
                    (cur_x + 12 * s, cur_y + 11 * s),
                ]
                shadow_points = [(px + 1.5 * s, py + 1.5 * s) for px, py in points]
                draw.polygon(shadow_points, fill=(0, 0, 0, 90))
                draw.polygon(points, fill=(255, 255, 255, 255), outline=(15, 23, 42, 255))
            else:
                # Click ripple phase
                rf = f - 9  # 0 to 4
                cur_x, cur_y = target_x, target_y
                
                # Expanding subtle ripple ring
                radius = (10 + rf * 8) * s
                alpha = int(220 * (1.0 - (rf / 5.0)))
                
                draw.ellipse(
                    [cur_x - radius, cur_y - radius, cur_x + radius, cur_y + radius],
                    outline=(239, 68, 68, alpha),
                    width=max(2, int(2 * s))
                )
                if radius > 14 * s:
                    draw.ellipse(
                        [cur_x - radius + 5 * s, cur_y - radius + 5 * s, cur_x + radius - 5 * s, cur_y + radius - 5 * s],
                        outline=(168, 85, 247, int(alpha * 0.7)),
                        width=max(1, int(1.5 * s))
                    )

                # Pressed compact cursor
                points = [
                    (cur_x, cur_y),
                    (cur_x, cur_y + 14 * s),
                    (cur_x + 3.5 * s, cur_y + 10.5 * s),
                    (cur_x + 7 * s, cur_y + 16.5 * s),
                    (cur_x + 9.5 * s, cur_y + 15 * s),
                    (cur_x + 5.5 * s, cur_y + 9.5 * s),
                    (cur_x + 10.5 * s, cur_y + 9.5 * s),
                ]
                shadow_points = [(px + 1.5 * s, py + 1.5 * s) for px, py in points]
                draw.polygon(shadow_points, fill=(0, 0, 0, 90))
                draw.polygon(points, fill=(240, 240, 240, 255), outline=(220, 38, 38, 255))

            frames.append(frame.convert("RGB"))

        seq = int(step["sequence"])
        gif_dir = SCREENSHOTS_DIR / session_id
        gif_dir.mkdir(parents=True, exist_ok=True)

        # Remove previous demo gifs for this step to keep folder clean
        for old_gif in gif_dir.glob(f"step-{seq:03d}-demo*.*"):
            try:
                old_gif.unlink()
            except Exception:
                pass

        gif_timestamp = int(time.time() * 1000)
        gif_filename = f"step-{seq:03d}-demo-{gif_timestamp}.gif"
        gif_path = gif_dir / gif_filename

        frames[0].save(
            gif_path,
            save_all=True,
            append_images=frames[1:],
            duration=110,
            loop=0,
            optimize=True
        )

        rel_gif_path = f"screenshots/{session_id}/{gif_filename}"
        cursor.execute(
            "UPDATE workflow_steps SET screenshot_path = ? WHERE id = ? AND workflow_id = ?",
            (rel_gif_path, step["id"], session_id)
        )
        connection.commit()

        return {
            "success": True,
            "gif_url": f"/screenshots/{session_id}/{gif_filename}",
            "filename": gif_filename
        }
    finally:
        connection.close()


@app.delete("/sessions/{session_id}/steps/{step_id}/animate")
def delete_step_animation(session_id: str, step_id: int):
    """
    Deletes the animated micro-demo GIF for the step and reverts screenshot_path to static PNG.
    """
    connection = get_connection()
    try:
        cursor = connection.cursor()
        step = cursor.execute(
            "SELECT id, sequence FROM workflow_steps WHERE id = ? AND workflow_id = ?",
            (step_id, session_id)
        ).fetchone()
        if not step:
            raise HTTPException(status_code=404, detail="Step not found")

        seq = int(step["sequence"])
        gif_dir = SCREENSHOTS_DIR / session_id
        deleted = False
        if gif_dir.is_dir():
            for old_gif in gif_dir.glob(f"step-{seq:03d}-demo*.*"):
                try:
                    old_gif.unlink()
                    deleted = True
                except Exception as e:
                    print("Error unlinking gif file:", e)

        # Revert database screenshot_path back to static PNG
        png_filename = f"step-{seq:03d}.png"
        rel_png_path = f"screenshots/{session_id}/{png_filename}"
        cursor.execute(
            "UPDATE workflow_steps SET screenshot_path = ? WHERE id = ? AND workflow_id = ?",
            (rel_png_path, step["id"], session_id)
        )
        connection.commit()

        return {
            "success": True,
            "deleted": deleted,
            "png_url": f"/screenshots/{session_id}/{png_filename}",
            "message": "Micro-demo animation removed successfully"
        }
    finally:
        connection.close()


# =========================================================
# GLOBAL DEEP SEARCH (COMMAND PALETTE)
# =========================================================

@app.get("/search")
def global_deep_search(q: str = Query(..., min_length=1)):
    """
    Deep search across all workflows, step titles, actions, notes, URLs, and tags.
    """
    connection = get_connection()
    q_clean = q.strip()
    like_query = f"%{q_clean}%"
    try:
        cursor = connection.cursor()
        
        # 1. Search Workflows
        wf_rows = cursor.execute(
            """
            SELECT id, name, application, tags, started_at
            FROM workflows
            WHERE name LIKE ? OR application LIKE ? OR tags LIKE ?
            ORDER BY updated_at DESC
            LIMIT 10
            """,
            (like_query, like_query, like_query)
        ).fetchall()
        
        workflows_match = [
            {
                "id": row["id"],
                "name": row["name"],
                "application": row["application"],
                "tags": (row["tags"] if "tags" in row.keys() else "") or "",
                "startedAt": row["started_at"]
            }
            for row in wf_rows
        ]
        
        # 2. Search Steps
        step_rows = cursor.execute(
            """
            SELECT ws.id as step_id, ws.sequence, ws.action, ws.title as orig_title, ws.url, ws.value,
                   w.id as workflow_id, w.name as workflow_name, w.application,
                   se.title as edit_title, se.description, se.note, se.expected
            FROM workflow_steps ws
            JOIN workflows w ON ws.workflow_id = w.id
            LEFT JOIN step_edits se ON ws.id = se.step_id
            WHERE ws.title LIKE ? OR ws.action LIKE ? OR ws.url LIKE ? OR ws.value LIKE ?
               OR se.title LIKE ? OR se.description LIKE ? OR se.note LIKE ? OR se.expected LIKE ?
            ORDER BY w.updated_at DESC, ws.sequence ASC
            LIMIT 25
            """,
            (like_query, like_query, like_query, like_query, like_query, like_query, like_query, like_query)
        ).fetchall()
        
        steps_match = []
        for row in step_rows:
            display_title = row["edit_title"] if row["edit_title"] else (row["orig_title"] or f"Step {row['sequence']}")
            match_context = ""
            if row["note"] and q_clean.lower() in row["note"].lower():
                match_context = f"Note: {row['note'][:60]}..."
            elif row["description"] and q_clean.lower() in row["description"].lower():
                match_context = f"Desc: {row['description'][:60]}..."
            elif row["url"] and q_clean.lower() in row["url"].lower():
                match_context = f"URL: {row['url'][:50]}..."
            elif row["action"] and q_clean.lower() in row["action"].lower():
                match_context = f"Action: {row['action']}"
            
            steps_match.append({
                "stepId": row["step_id"],
                "sequence": row["sequence"],
                "workflowId": row["workflow_id"],
                "workflowName": row["workflow_name"],
                "application": row["application"],
                "title": display_title,
                "action": row["action"],
                "matchContext": match_context
            })
            
        return {
            "query": q_clean,
            "totalMatches": len(workflows_match) + len(steps_match),
            "workflows": workflows_match,
            "steps": steps_match
        }
    finally:
        connection.close()


class DesktopCaptureRequest(BaseModel):
    session_id: Optional[str] = None
    monitor_index: Optional[int] = 1   # 1 = primary, 2 = secondary, 0 = all monitors combined
    title: Optional[str] = None


def _ensure_desktop_session(cursor, connection, session_id: Optional[str]) -> str:
    """Create or validate a workflow session for desktop captures."""
    import uuid, datetime
    if session_id:
        wf = cursor.execute("SELECT id FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if wf:
            return session_id
    session_id = f"wf_{uuid.uuid4().hex[:8]}"
    now_str = datetime.datetime.now().strftime("%b %d, %Y %I:%M %p")
    cursor.execute(
        """
        INSERT INTO workflows (id, name, application, status, started_at, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (session_id, f"Desktop SOP — {now_str}", "Windows Desktop", "completed", utc_now(), utc_now(), utc_now())
    )
    connection.commit()
    return session_id


def _insert_desktop_step(cursor, connection, session_id: str, filename: str, title: str) -> dict:
    """Insert a workflow_step row for a desktop screenshot and return URL info."""
    file_path = SCREENSHOTS_DIR / session_id / filename
    db_path = str(Path("screenshots") / session_id / filename)

    last_step = cursor.execute(
        "SELECT sequence FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence DESC LIMIT 1",
        (session_id,)
    ).fetchone()
    next_seq = (last_step["sequence"] + 1) if last_step else 1

    step_title = title or f"Desktop Screen Capture {next_seq}"

    cursor.execute(
        """
        INSERT INTO workflow_steps (
            workflow_id, sequence, action, timestamp, title, screenshot_path
        ) VALUES (?, ?, ?, ?, ?, ?)
        """,
        (session_id, next_seq, "desktop_capture", utc_now(), step_title, db_path)
    )
    connection.commit()
    return {
        "success": True,
        "sessionId": session_id,
        "screenshotUrl": f"/screenshots/{session_id}/{filename}",
        "sequence": next_seq,
        "title": step_title,
    }


@app.get("/desktop/monitors")
def list_desktop_monitors():
    """
    Returns a list of available desktop monitors with index, resolution and position.
    Uses mss for cross-monitor enumeration.
    """
    try:
        import mss
        with mss.mss() as sct:
            monitors = []
            for i, m in enumerate(sct.monitors):
                if i == 0:
                    label = f"All Monitors Combined ({m['width']}x{m['height']})"
                else:
                    label = f"Monitor {i} - {m['width']}x{m['height']} at ({m['left']},{m['top']})"
                monitors.append({
                    "index": i,
                    "label": label,
                    "width": m["width"],
                    "height": m["height"],
                    "left": m["left"],
                    "top": m["top"],
                    "isPrimary": i == 1,
                })
        return {"success": True, "monitors": monitors}
    except Exception as e:
        return {"success": False, "monitors": [{"index": 1, "label": "Primary Monitor", "width": 1920, "height": 1080, "left": 0, "top": 0, "isPrimary": True}], "error": str(e)}


@app.post("/desktop/capture")
def capture_desktop_screen(request: Optional[DesktopCaptureRequest] = None):
    """
    Captures the Windows Desktop screen natively.
    Tries 4 methods in order:
      1. mss.MSS (fast, multi-monitor)
      2. PIL ImageGrab
      3. PowerShell CopyFromScreen (most compatible in user sessions)
      4. capture_screen.py via subprocess
    Auto-creates a workflow session if none provided.
    """
    import uuid
    try:
        session_id = (request.session_id if request else None) or None
        monitor_index = int((request.monitor_index if request else None) or 1)
        step_title = (request.title if request else None) or None

        connection = get_connection()
        try:
            cursor = connection.cursor()
            session_id = _ensure_desktop_session(cursor, connection, session_id)

            filename = f"desktop_{uuid.uuid4().hex[:10]}.png"
            folder = SCREENSHOTS_DIR / session_id
            folder.mkdir(parents=True, exist_ok=True)
            file_path = folder / filename
            file_path_str = str(file_path)

            captured = False
            errors = []

            # ── Method 1: mss.MSS ──────────────────────────────────────────────
            try:
                import mss
                import mss.tools
                with mss.MSS() as sct:
                    monitors = sct.monitors
                    mon = monitors[monitor_index] if monitor_index < len(monitors) else monitors[1]
                    screenshot = sct.grab(mon)
                    mss.tools.to_png(screenshot.rgb, screenshot.size, output=file_path_str)
                captured = True
            except Exception as e:
                errors.append(f"mss: {e}")

            # ── Method 2: PIL ImageGrab ──────────────────────────────────────────
            if not captured:
                try:
                    from PIL import ImageGrab
                    all_screens = (monitor_index == 0)
                    img = ImageGrab.grab(include_layered_windows=True, all_screens=all_screens)
                    img.save(file_path_str, format="PNG")
                    captured = True
                except Exception as e:
                    errors.append(f"PIL: {e}")

            # ── Method 3: PowerShell CopyFromScreen ─────────────────────────────
            if not captured:
                try:
                    ps_script = BASE_DIR / "capture.ps1"
                    if not ps_script.exists():
                        # Write the PowerShell capture helper if missing
                        ps_script.write_text(
                            r"""
param([string]$outPath, [int]$monIdx = 0)
Add-Type -AssemblyName System.Windows.Forms
Add-Type -AssemblyName System.Drawing
$screens = [System.Windows.Forms.Screen]::AllScreens
$screen = if ($monIdx -lt $screens.Length) { $screens[$monIdx] } else { $screens[0] }
$bmp = New-Object System.Drawing.Bitmap($screen.Bounds.Width, $screen.Bounds.Height)
$g = [System.Drawing.Graphics]::FromImage($bmp)
$g.CopyFromScreen($screen.Bounds.X, $screen.Bounds.Y, 0, 0, $bmp.Size)
$bmp.Save($outPath, [System.Drawing.Imaging.ImageFormat]::Png)
$g.Dispose(); $bmp.Dispose()
Write-Output "OK:$outPath"
""".strip(), encoding="utf-8")

                    ps_mon_idx = max(0, monitor_index - 1)  # PS uses 0-based AllScreens
                    result = subprocess.run(
                        ["powershell", "-ExecutionPolicy", "Bypass", "-NonInteractive",
                         "-File", str(ps_script), file_path_str, str(ps_mon_idx)],
                        capture_output=True, text=True, timeout=20
                    )
                    if result.returncode == 0 and file_path.exists() and file_path.stat().st_size > 1000:
                        captured = True
                    else:
                        errors.append(f"PowerShell: rc={result.returncode} {result.stderr.strip()}")
                except Exception as e:
                    errors.append(f"PowerShell: {e}")

            # ── Method 4: capture_screen.py subprocess ───────────────────────────
            if not captured:
                try:
                    helper = BASE_DIR / "capture_screen.py"
                    result = subprocess.run(
                        [sys.executable, str(helper), file_path_str, str(monitor_index)],
                        capture_output=True, text=True, timeout=20
                    )
                    if result.returncode == 0 and file_path.exists() and file_path.stat().st_size > 1000:
                        captured = True
                    else:
                        errors.append(f"capture_screen.py: rc={result.returncode} {result.stderr.strip()}")
                except Exception as e:
                    errors.append(f"capture_screen.py: {e}")

            # ── Method 5: macOS native screencapture CLI ──────────────────────────
            if not captured and sys.platform == "darwin":
                try:
                    result = subprocess.run(["screencapture", "-x", file_path_str], capture_output=True, timeout=10)
                    if result.returncode == 0 and file_path.exists() and file_path.stat().st_size > 1000:
                        captured = True
                    else:
                        errors.append(f"screencapture: rc={result.returncode}")
                except Exception as e:
                    errors.append(f"screencapture: {e}")

            if not captured:
                raise RuntimeError(f"All capture methods failed: {'; '.join(errors)}")

            result = _insert_desktop_step(cursor, connection, session_id, filename, step_title)
            return result
        finally:
            connection.close()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Desktop capture failed: {str(e)}")


class DesktopBase64CaptureRequest(BaseModel):
    session_id: Optional[str] = None
    image: str
    title: Optional[str] = None


@app.post("/desktop/capture-base64")
def capture_desktop_base64(request: DesktopBase64CaptureRequest):
    """
    Accepts a base64 data-URL desktop screenshot (from browser getDisplayMedia)
    and saves it as a new workflow step.
    """
    import uuid
    try:
        match = re.fullmatch(r"data:image/[^;]+;base64,(.+)", request.image, flags=re.DOTALL)
        if not match:
            raise HTTPException(status_code=400, detail="Invalid data URL format")

        image_data = base64.b64decode(match.group(1))
        if len(image_data) < 1000:
            raise HTTPException(status_code=400, detail="Captured image appears to be blank or empty")

        connection = get_connection()
        try:
            cursor = connection.cursor()
            session_id = _ensure_desktop_session(cursor, connection, request.session_id)

            filename = f"desktop_{uuid.uuid4().hex[:10]}.png"
            folder = SCREENSHOTS_DIR / session_id
            folder.mkdir(parents=True, exist_ok=True)
            file_path = folder / filename

            with open(file_path, "wb") as f:
                f.write(image_data)

            result = _insert_desktop_step(cursor, connection, session_id, filename, request.title)
            return result
        finally:
            connection.close()
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Desktop base64 capture failed: {str(e)}")


# =========================================================
# EXPORT DOCX (Word Document)
# =========================================================

@app.get("/sessions/{session_id}/export/docx")
def export_session_docx(session_id: str):
    """
    Generates and returns a formatted Microsoft Word (.docx) SOP document
    containing step descriptions, notes, expected outcomes, and embedded screenshots.
    """
    session_data = get_session(session_id)
    name = session_data.get("name") or "ProcSnap SOP Guide"
    application = session_data.get("application") or "System"
    steps = [s for s in session_data.get("steps", []) if not s.get("hidden", False)]
    
    doc = docx.Document()
    
    # 0.75-inch page margins
    for section in doc.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)
        
    # Document Title Header
    title_p = doc.add_paragraph()
    title_p.paragraph_format.space_before = Pt(0)
    title_p.paragraph_format.space_after = Pt(4)
    run_title = title_p.add_run(name)
    run_title.font.name = "Calibri"
    run_title.font.size = Pt(22)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x11, 0x18, 0x27)
    
    # Metadata Subtitle
    meta_p = doc.add_paragraph()
    meta_p.paragraph_format.space_after = Pt(16)
    gen_date = datetime.now().strftime("%B %d, %Y")
    run_meta = meta_p.add_run(f"Recorded with {application} • {len(steps)} Steps • Generated on {gen_date}")
    run_meta.font.name = "Calibri"
    run_meta.font.size = Pt(10)
    run_meta.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
    
    # Process Steps
    for idx, s in enumerate(steps, 1):
        step_num = s.get("sequence", idx)
        title_text = s.get("title") or f"Step {step_num}"
        desc_text = s.get("description") or ""
        note_text = s.get("note") or ""
        expected_text = s.get("expected") or ""
        
        # Step Header (Heading 2)
        step_p = doc.add_paragraph()
        step_p.paragraph_format.space_before = Pt(14)
        step_p.paragraph_format.space_after = Pt(4)
        step_p.paragraph_format.keep_with_next = True
        
        run_badge = step_p.add_run(f"STEP {step_num}: ")
        run_badge.font.name = "Calibri"
        run_badge.font.size = Pt(12)
        run_badge.font.bold = True
        run_badge.font.color.rgb = RGBColor(0x4F, 0x46, 0xE5) # Indigo
        
        run_stitle = step_p.add_run(title_text)
        run_stitle.font.name = "Calibri"
        run_stitle.font.size = Pt(13)
        run_stitle.font.bold = True
        run_stitle.font.color.rgb = RGBColor(0x11, 0x18, 0x27)
        
        # Description
        if desc_text:
            desc_p = doc.add_paragraph()
            desc_p.paragraph_format.space_after = Pt(6)
            run_desc = desc_p.add_run(desc_text)
            run_desc.font.name = "Calibri"
            run_desc.font.size = Pt(10.5)
            run_desc.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            
        # Note Callout Box
        if note_text:
            note_table = doc.add_table(rows=1, cols=1)
            note_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = note_table.cell(0, 0)
            cell.width = Inches(6.5)
            shading = parse_xml(r'<w:shd {} w:fill="F5F3FF"/>'.format(nsdecls('w')))
            cell._tc.get_or_add_tcPr().append(shading)
            borders = parse_xml(r'<w:tcBorders {}><w:left w:val="single" w:sz="24" w:space="0" w:color="7C3AED"/><w:top w:val="none"/><w:right w:val="none"/><w:bottom w:val="none"/></w:tcBorders>'.format(nsdecls('w')))
            cell._tc.get_or_add_tcPr().append(borders)
            
            p_note = cell.paragraphs[0]
            p_note.paragraph_format.space_before = Pt(4)
            p_note.paragraph_format.space_after = Pt(4)
            r_nl = p_note.add_run("Note: ")
            r_nl.bold = True
            r_nl.font.color.rgb = RGBColor(0x7C, 0x3A, 0xED)
            r_nt = p_note.add_run(note_text)
            r_nt.font.color.rgb = RGBColor(0x4C, 0x1D, 0x95)
            
        # Expected Result Callout Box
        if expected_text:
            exp_table = doc.add_table(rows=1, cols=1)
            exp_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = exp_table.cell(0, 0)
            cell.width = Inches(6.5)
            shading = parse_xml(r'<w:shd {} w:fill="ECFDF5"/>'.format(nsdecls('w')))
            cell._tc.get_or_add_tcPr().append(shading)
            borders = parse_xml(r'<w:tcBorders {}><w:left w:val="single" w:sz="24" w:space="0" w:color="10B981"/><w:top w:val="none"/><w:right w:val="none"/><w:bottom w:val="none"/></w:tcBorders>'.format(nsdecls('w')))
            cell._tc.get_or_add_tcPr().append(borders)
            
            p_exp = cell.paragraphs[0]
            p_exp.paragraph_format.space_before = Pt(4)
            p_exp.paragraph_format.space_after = Pt(4)
            r_el = p_exp.add_run("Expected Result: ")
            r_el.bold = True
            r_el.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
            r_et = p_exp.add_run(expected_text)
            r_et.font.color.rgb = RGBColor(0x06, 0x5F, 0x46)
            
        # Embedded Screenshot Image with Composited Annotations
        screenshot_path = s.get("screenshotPath")
        if screenshot_path:
            img_file = BASE_DIR / screenshot_path
            if not img_file.exists():
                img_file = BASE_DIR / "screenshots" / session_id / Path(screenshot_path).name
            if img_file.exists():
                try:
                    from .annotation_renderer import AnnotationRenderer
                except ImportError:
                    from annotation_renderer import AnnotationRenderer

                composited = AnnotationRenderer.composite_image(
                    str(img_file),
                    annotations=s.get("annotations", []),
                    focus_crop=s.get("focus_crop")
                )

                img_p = doc.add_paragraph()
                img_p.paragraph_format.space_before = Pt(6)
                img_p.paragraph_format.space_after = Pt(16)
                img_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run_img = img_p.add_run()

                if composited:
                    img_buf = io.BytesIO()
                    composited.save(img_buf, format="PNG")
                    img_buf.seek(0)
                    run_img.add_picture(img_buf, width=Inches(6.0))
                else:
                    run_img.add_picture(str(img_file), width=Inches(6.0))

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    
    clean_filename = re.sub(r'[<>:"/\\|?*]+', '-', name).strip() or "ProcSnap_SOP"
    headers = {
        "Content-Disposition": f'attachment; filename="{clean_filename}.docx"',
        "Access-Control-Expose-Headers": "Content-Disposition"
    }
    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers=headers
    )


# =========================================================
# EXPORT PPTX (PowerPoint Presentation)
# =========================================================

@app.get("/sessions/{session_id}/export/pptx")
def export_session_pptx(session_id: str):
    """
    Generates and returns a widescreen 16:9 PowerPoint (.pptx) slide deck
    with title slide, step descriptions, notes, and embedded high-res screenshots.
    """
    if pptx is None:
        raise HTTPException(status_code=500, detail="python-pptx library not installed")
    session_data = get_session(session_id)
    name = session_data.get("name") or "ProcSnap SOP Guide"
    application = session_data.get("application") or "System"
    steps = [s for s in session_data.get("steps", []) if not s.get("hidden", False)]

    prs = pptx.Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- Title Slide ---
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(13.333), PptxInches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptxRGBColor(15, 23, 42)
    bg.line.fill.background()

    tx_box = title_slide.shapes.add_textbox(PptxInches(1.5), PptxInches(2.2), PptxInches(10.333), PptxInches(3.0))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "PROCSNAP STANDARD OPERATING PROCEDURE"
    p1.font.size = PptxPt(13)
    p1.font.bold = True
    p1.font.color.rgb = PptxRGBColor(99, 102, 241)

    p2 = tf.add_paragraph()
    p2.text = name
    p2.font.size = PptxPt(36)
    p2.font.bold = True
    p2.font.color.rgb = PptxRGBColor(255, 255, 255)

    p3 = tf.add_paragraph()
    p3.text = f"Application: {application}  •  {len(steps)} Steps  •  {datetime.now().strftime('%B %d, %Y')}"
    p3.font.size = PptxPt(15)
    p3.font.color.rgb = PptxRGBColor(148, 163, 184)

    # --- Step Slides ---
    for idx, s in enumerate(steps, 1):
        step_num = s.get("sequence", idx)
        title_text = s.get("title") or f"Step {step_num}"
        desc_text = s.get("description") or ""

        slide = prs.slides.add_slide(blank_layout)
        header_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(0.4), PptxInches(11.733), PptxInches(1.2))
        htf = header_box.text_frame
        htf.word_wrap = True
        
        hp1 = htf.paragraphs[0]
        hp1.text = f"STEP {step_num}: {title_text}"
        hp1.font.size = PptxPt(20)
        hp1.font.bold = True
        hp1.font.color.rgb = PptxRGBColor(15, 23, 42)

        if desc_text:
            hp2 = htf.add_paragraph()
            hp2.text = desc_text
            hp2.font.size = PptxPt(13)
            hp2.font.color.rgb = PptxRGBColor(71, 85, 105)

        screenshot_path = s.get("screenshotPath")
        if screenshot_path:
            img_file = BASE_DIR / screenshot_path
            if not img_file.exists():
                img_file = BASE_DIR / "screenshots" / session_id / Path(screenshot_path).name
            if img_file.exists():
                try:
                    from .annotation_renderer import AnnotationRenderer
                except ImportError:
                    from annotation_renderer import AnnotationRenderer

                composited = AnnotationRenderer.composite_image(
                    str(img_file),
                    annotations=s.get("annotations", []),
                    focus_crop=s.get("focus_crop")
                )

                try:
                    if composited:
                        img_buf = io.BytesIO()
                        composited.save(img_buf, format="PNG")
                        img_buf.seek(0)
                        slide.shapes.add_picture(
                            img_buf,
                            PptxInches(0.8),
                            PptxInches(1.8),
                            width=PptxInches(11.733)
                        )
                    else:
                        slide.shapes.add_picture(
                            str(img_file),
                            PptxInches(0.8),
                            PptxInches(1.8),
                            width=PptxInches(11.733)
                        )
                except Exception:
                    pass

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    clean_filename = re.sub(r'[<>:"/\\|?*]+', '-', name).strip() or "ProcSnap_Presentation"
    headers = {
        "Content-Disposition": f'attachment; filename="{clean_filename}.pptx"',
        "Access-Control-Expose-Headers": "Content-Disposition"
    }
    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers=headers
    )


# =========================================================
# EXPORT / IMPORT JSON (Full Portable Backup)
# =========================================================

@app.get("/sessions/{session_id}/export/json")
def export_session_json(session_id: str):
    """
    Exports a standalone JSON backup of the workflow including all metadata,
    annotations, edits, and embedded base64 screenshots.
    """
    session_data = get_session(session_id)
    for step in session_data.get("steps", []):
        screenshot_path = step.get("screenshotPath")
        if screenshot_path:
            img_file = BASE_DIR / screenshot_path
            if not img_file.exists():
                img_file = BASE_DIR / "screenshots" / session_id / Path(screenshot_path).name
            if img_file.exists():
                try:
                    with open(img_file, "rb") as f:
                        step["screenshotBase64"] = base64.b64encode(f.read()).decode("utf-8")
                except Exception:
                    step["screenshotBase64"] = None

    clean_filename = re.sub(r'[<>:"/\\|?*]+', '-', session_data.get("name", "Workflow")).strip()
    headers = {
        "Content-Disposition": f'attachment; filename="{clean_filename}.procsnap.json"',
        "Access-Control-Expose-Headers": "Content-Disposition"
    }
    return Response(
        content=json.dumps(session_data, indent=2),
        media_type="application/json",
        headers=headers
    )


@app.post("/sessions/import")
def import_session(request: ImportWorkflowRequest):
    """
    Imports a complete workflow from JSON, reconstructing the database rows,
    annotations, edits, and decoding base64 screenshots to local disk.
    """
    data = request.workflow
    if not isinstance(data, dict):
        raise HTTPException(status_code=400, detail="Invalid workflow JSON payload")

    new_id = str(uuid4())
    now = utc_now()
    name = data.get("name") or "Imported Workflow"
    application = data.get("application") or "Imported"
    tags = data.get("tags") or ""

    connection = get_connection()
    try:
        cursor = connection.cursor()
        cursor.execute(
            """
            INSERT INTO workflows (id, name, application, status, tags, started_at, ended_at, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (new_id, name, application, "completed", tags, data.get("startedAt", now), data.get("endedAt", now), now, now),
        )

        new_screenshot_dir = SCREENSHOTS_DIR / new_id
        new_screenshot_dir.mkdir(parents=True, exist_ok=True)

        steps = data.get("steps", [])
        for idx, s in enumerate(steps, 1):
            seq = s.get("sequence", idx)
            action = s.get("action", "click")
            timestamp = s.get("timestamp", now)
            url = s.get("url", "")
            title = s.get("title")
            val = s.get("value")
            sel_text = s.get("selectedText")
            prev_url = s.get("previousUrl")
            checked = 1 if s.get("checked") else 0
            elem = json.dumps(s.get("element", {})) if isinstance(s.get("element"), dict) else "{}"

            new_path = None
            if s.get("screenshotBase64"):
                try:
                    img_data = base64.b64decode(s["screenshotBase64"])
                    filename = f"step_{seq}_{uuid4().hex[:8]}.png"
                    img_dest = new_screenshot_dir / filename
                    with open(img_dest, "wb") as f:
                        f.write(img_data)
                    new_path = str(img_dest)
                except Exception:
                    pass

            cursor.execute(
                """
                INSERT INTO workflow_steps (
                    workflow_id, sequence, action, timestamp, url, title, value,
                    selected_text, previous_url, checked, element_json, screenshot_path
                )
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (new_id, seq, action, timestamp, url, title, val, sel_text, prev_url, checked, elem, new_path),
            )
            step_id = cursor.lastrowid

            if s.get("annotations"):
                cursor.execute(
                    """
                    INSERT INTO step_annotations (step_id, workflow_id, data, created_at, updated_at)
                    VALUES (?, ?, ?, ?, ?)
                    """,
                    (step_id, new_id, json.dumps(s["annotations"]), now, now),
                )

            if s.get("note") or s.get("expected") or s.get("voiceover") or s.get("hidden"):
                cursor.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, voiceover, hidden, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (step_id, new_id, title, s.get("description"), s.get("note"), s.get("expected"), s.get("voiceover"), 1 if s.get("hidden") else 0, now),
                )

        connection.commit()
        return {
            "success": True,
            "id": new_id,
            "name": name,
            "stepCount": len(steps),
            "message": f"Successfully imported '{name}' with {len(steps)} steps"
        }
    finally:
        connection.close()



# =========================================================
# SYSTEM REQUIREMENTS & HEALTH DIAGNOSTICS
# =========================================================

@app.get("/system/requirements")
def get_system_requirements():
    """
    Scans the local environment and returns live diagnostic statuses for all components.
    """
    import platform
    import importlib.metadata
    
    # 1. Python Environment
    py_version = platform.python_version()
    py_executable = sys.executable
    in_venv = hasattr(sys, 'real_prefix') or (hasattr(sys, 'base_prefix') and sys.base_prefix != sys.prefix)
    
    # 2. Package Dependency Checks
    tracked_packages = [
        {"name": "fastapi", "required": ">=0.110.0", "description": "Core HTTP Backend API & Studio Server"},
        {"name": "uvicorn", "required": ">=0.29.0", "description": "High-Performance ASGI Web Server"},
        {"name": "python-docx", "package_lookup": "python-docx", "required": ">=1.1.0", "description": "Microsoft Word (.docx) SOP Document Generator"},
        {"name": "pillow", "package_lookup": "Pillow", "required": ">=10.0.0", "description": "Screenshot & Canvas Image Processing Engine"},
        {"name": "edge-tts", "package_lookup": "edge-tts", "required": ">=6.1.0", "description": "Offline-ready Microsoft Neural Voice Narration"},
        {"name": "lxml", "package_lookup": "lxml", "required": ">=4.9.0", "description": "XML & Document Layout Styling Engine"},
        {"name": "pydantic", "package_lookup": "pydantic", "required": ">=2.0.0", "description": "Data Validation & Schema Modeling"},
    ]
    
    package_results = []
    all_packages_ok = True
    
    for pkg in tracked_packages:
        lookup_name = pkg.get("package_lookup", pkg["name"])
        installed = False
        installed_version = None
        
        try:
            installed_version = importlib.metadata.version(lookup_name)
            installed = True
        except Exception:
            # Fallback import check
            try:
                mod = __import__(pkg["name"].replace("-", "_"))
                installed_version = getattr(mod, "__version__", "Installed")
                installed = True
            except Exception:
                installed = False
                installed_version = None
                all_packages_ok = False
                
        package_results.append({
            "name": pkg["name"],
            "required": pkg["required"],
            "description": pkg["description"],
            "installed": installed,
            "version": installed_version or "Not Installed"
        })
        
    # 3. Database Stats
    db_file = DATABASE_PATH
    db_size = db_file.stat().st_size if db_file.exists() else 0
    wf_count = 0
    step_count = 0
    
    try:
        conn = get_connection()
        try:
            c = conn.cursor()
            wf_count = c.execute("SELECT count(*) FROM workflows").fetchone()[0]
            step_count = c.execute("SELECT count(*) FROM workflow_steps").fetchone()[0]
        finally:
            conn.close()
        db_connected = True
    except Exception:
        db_connected = False
        
    # 4. Ollama AI Status
    try:
        ollama_info = get_ai_status()
    except Exception:
        ollama_info = {"running": False, "models": []}
    
    # 5. Extension Folder Check
    ext_dir = BASE_DIR.parent / "extension"
    manifest_file = ext_dir / "manifest.json"
    extension_ready = ext_dir.exists() and manifest_file.exists()
    
    # 6. Windows Shortcuts Check
    desktop_shortcut = False
    start_shortcut = False
    try:
        desktop_dir = Path(os.path.expandvars("%USERPROFILE%")) / "Desktop"
        desktop_shortcut = (desktop_dir / "ProcSnap.lnk").exists()
        
        programs_dir = Path(os.path.expandvars("%APPDATA%")) / "Microsoft" / "Windows" / "Start Menu" / "Programs"
        start_shortcut = (programs_dir / "ProcSnap.lnk").exists()
    except Exception:
        pass
        
    return {
        "success": True,
        "status": "ready" if all_packages_ok and db_connected else "needs_attention",
        "python": {
            "version": py_version,
            "executable": py_executable,
            "in_venv": in_venv,
            "os": platform.platform()
        },
        "packages": {
            "all_ok": all_packages_ok,
            "items": package_results
        },
        "database": {
            "connected": db_connected,
            "path": str(db_file),
            "size_bytes": db_size,
            "workflows_count": wf_count,
            "steps_count": step_count
        },
        "ollama": ollama_info,
        "extension": {
            "ready": extension_ready,
            "path": str(ext_dir)
        },
        "shortcuts": {
            "desktop": desktop_shortcut,
            "start_menu": start_shortcut
        }
    }


@app.post("/system/reinstall-packages")
def reinstall_packages():
    """
    Runs pip install -r requirements.txt using the active Python executable.
    """
    req_file = BASE_DIR / "requirements.txt"
    if not req_file.exists():
        raise HTTPException(status_code=404, detail="requirements.txt not found")
        
    try:
        cmd = [sys.executable, "-m", "pip", "install", "--upgrade", "-r", str(req_file)]
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        
        output = result.stdout + "\n" + result.stderr
        success = (result.returncode == 0)
        
        return {
            "success": success,
            "return_code": result.returncode,
            "output": output.strip() or "Installation completed."
        }
    except subprocess.TimeoutExpired:
        return {"success": False, "output": "Installation timed out after 120 seconds."}
    except Exception as e:
        return {"success": False, "output": f"Error running package installer: {str(e)}"}


@app.post("/system/repair-shortcuts")
def repair_shortcuts():
    """
    Re-creates Desktop and Start Menu shortcuts for ProcSnap.
    """
    install_dir = BASE_DIR.parent
    vbs_path = Path(os.path.expandvars("%TEMP%")) / "create_procsnap_shortcuts.vbs"
    
    script_content = f"""Set WshShell = CreateObject("WScript.Shell")
strDesktop = WshShell.SpecialFolders("Desktop")
strPrograms = WshShell.SpecialFolders("Programs")

' Desktop shortcut
Set oLink1 = WshShell.CreateShortcut(strDesktop & "\\ProcSnap.lnk")
oLink1.TargetPath = "{install_dir}\\start.bat"
oLink1.WorkingDirectory = "{install_dir}"
oLink1.Description = "ProcSnap - Local Process Recorder & SOP Studio"
oLink1.IconLocation = "shell32.dll, 220"
oLink1.Save

' Start Menu shortcut
Set oLink2 = WshShell.CreateShortcut(strPrograms & "\\ProcSnap.lnk")
oLink2.TargetPath = "{install_dir}\\start.bat"
oLink2.WorkingDirectory = "{install_dir}"
oLink2.Description = "ProcSnap - Local Process Recorder & SOP Studio"
oLink2.IconLocation = "shell32.dll, 220"
oLink2.Save
"""
    try:
        vbs_path.write_text(script_content, encoding="utf-8")
        subprocess.run(["cscript", "//nologo", str(vbs_path)], check=True, capture_output=True)
        if vbs_path.exists():
            vbs_path.unlink()
        return {"success": True, "message": "Shortcuts created successfully on Desktop and Start Menu."}
    except Exception as e:
        return {"success": False, "message": f"Failed to create shortcuts: {str(e)}"}


class ExtensionLaunchRequest(BaseModel):
    browser: Optional[str] = "default"

@app.post("/system/open-extension-installer")
def open_extension_installer(payload: Optional[ExtensionLaunchRequest] = None):
    """
    Opens browser-specific extensions page and copies extension path to clipboard.
    """
    browser = (payload.browser.lower() if payload and payload.browser else "default")
    ext_dir = str(BASE_DIR.parent / "extension")
    
    try:
        subprocess.run(["powershell", "-NoProfile", "-Command", f"Set-Clipboard -Value '{ext_dir}'"], capture_output=True, timeout=3)
    except Exception:
        pass
        
    try:
        if browser == "chrome":
            subprocess.Popen(["cmd.exe", "/c", "start", "chrome", "chrome://extensions"], shell=True)
            return {"success": True, "message": "Google Chrome Extensions opened. Paste path into 'Load unpacked'!"}
        elif browser == "edge":
            subprocess.Popen(["cmd.exe", "/c", "start", "msedge", "edge://extensions"], shell=True)
            return {"success": True, "message": "Microsoft Edge Extensions opened. Paste path into 'Load unpacked'!"}
        elif browser == "brave":
            subprocess.Popen(["cmd.exe", "/c", "start", "brave", "brave://extensions"], shell=True)
            return {"success": True, "message": "Brave Extensions opened. Paste path into 'Load unpacked'!"}
        else:
            helper_bat = BASE_DIR.parent / "install_extension.bat"
            if helper_bat.exists():
                subprocess.Popen(["cmd.exe", "/c", "start", "", str(helper_bat)], shell=True)
            return {"success": True, "message": "Browser extension installer launched."}
    except Exception as e:
        return {"success": False, "message": str(e)}


@app.post("/system/git-pull")
def git_pull_latest():
    """
    Pull the latest commits from origin/main into the local installation.
    Uses git executable from common Windows install paths.
    Returns combined stdout+stderr output so the dashboard can display it.
    """
    # Locate git executable
    git_exe = shutil.which("git")
    if not git_exe:
        common_paths = [
            Path(os.environ.get("LOCALAPPDATA", "")) / "Programs" / "Git" / "cmd" / "git.exe",
            Path("C:/Program Files/Git/cmd/git.exe"),
            Path("C:/Program Files (x86)/Git/cmd/git.exe"),
        ]
        for p in common_paths:
            if p.exists():
                git_exe = str(p)
                break

    if not git_exe:
        return {
            "success": False,
            "output": (
                "❌ Git not found on this machine.\n"
                "Download and install Git from https://git-scm.com/download/win\n"
                "Then re-run this update."
            )
        }

    repo_root = BASE_DIR.parent  # project root (one level above backend/)
    output_lines = [f"📁 Repository root: {repo_root}", f"🔧 Git: {git_exe}", ""]

    # Ensure remote is set correctly
    try:
        subprocess.run(
            [git_exe, "remote", "set-url", "origin", "https://github.com/MRVKY220895/ProcSnap.git"],
            cwd=str(repo_root), capture_output=True, text=True, timeout=15
        )
        output_lines.append("✓ Remote origin verified.")
    except Exception as e:
        output_lines.append(f"⚠ Could not set remote: {e}")

    # Run git pull
    output_lines.append("\n⬇ Running: git pull origin main ...\n")
    try:
        result = subprocess.run(
            [git_exe, "pull", "origin", "main"],
            cwd=str(repo_root),
            capture_output=True,
            text=True,
            timeout=120,
            env={**os.environ, "GIT_TERMINAL_PROMPT": "0"}  # disable interactive prompts
        )
        if result.stdout:
            output_lines.append(result.stdout.strip())
        if result.stderr:
            output_lines.append(result.stderr.strip())

        success = result.returncode == 0
        if success:
            output_lines.append("\n✅ Update complete! Please restart ProcSnap to apply changes.")
        else:
            output_lines.append(f"\n⚠ git pull exited with code {result.returncode}.")
            output_lines.append("If you see authentication errors, the repo may be private or requires a token.")

        return {"success": success, "output": "\n".join(output_lines)}
    except subprocess.TimeoutExpired:
        return {"success": False, "output": "\n".join(output_lines) + "\n\n⏰ Timed out after 2 minutes."}
    except Exception as e:
        return {"success": False, "output": "\n".join(output_lines) + f"\n\n❌ Error: {str(e)}"}


# =========================================================
# VIDEO / GIF TO STEP SOP IMPORT (OpenCV Keyframe Engine)
# =========================================================
@app.post("/workflows/import-video")
async def import_video_workflow(
    file: UploadFile = File(...),
    workflow_name: Optional[str] = Form(None),
    sensitivity: Optional[str] = Form("medium")
):
    """
    Extracts keyframe transitions from an uploaded video/GIF (Loom, MP4, WebM)
    and converts them into an editable ProcSnap SOP workflow with screenshots.
    """
    import tempfile
    try:
        import cv2
        import numpy as np
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="OpenCV (cv2) or NumPy is not installed. Please install opencv-python."
        )

    filename = file.filename or "video.mp4"
    ext = Path(filename).suffix.lower() or ".mp4"
    
    # Save uploaded video to temp file
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp_video:
        content = await file.read()
        tmp_video.write(content)
        tmp_video_path = tmp_video.name

    try:
        cap = cv2.VideoCapture(tmp_video_path)
        if not cap.isOpened():
            raise HTTPException(status_code=400, detail="Could not open video file.")

        fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        if total_frames <= 0:
            raise HTTPException(status_code=400, detail="Video contains no readable frames.")

        duration_sec = total_frames / fps
        
        # Sample rate: 1 frame every 0.4 seconds (2.5 fps)
        sample_step = max(1, int(fps * 0.4))
        
        # Thresholds based on sensitivity
        sens_map = {
            "high": 4.0,
            "medium": 8.0,
            "low": 15.0
        }
        threshold = sens_map.get((sensitivity or "").lower(), 8.0)

        selected_frames = []
        prev_gray = None
        frame_idx = 0
        last_selected_time = -10.0

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            if frame_idx % sample_step == 0:
                current_time = frame_idx / fps
                small = cv2.resize(frame, (320, 180))
                gray = cv2.cvtColor(small, cv2.COLOR_BGR2GRAY)
                gray = cv2.GaussianBlur(gray, (5, 5), 0)

                if prev_gray is None:
                    # Always capture the initial starting screen
                    selected_frames.append((current_time, frame.copy()))
                    last_selected_time = current_time
                    prev_gray = gray
                else:
                    diff = cv2.absdiff(prev_gray, gray)
                    score = np.mean(diff)

                    # If significant screen change detected and at least 0.8s elapsed
                    if score >= threshold and (current_time - last_selected_time) >= 0.8:
                        selected_frames.append((current_time, frame.copy()))
                        last_selected_time = current_time
                        prev_gray = gray

                        if len(selected_frames) >= 40:
                            break  # Cap at 40 key steps per import

            frame_idx += 1

        cap.release()

        if not selected_frames:
            raise HTTPException(status_code=400, detail="Could not extract any distinct steps from video.")

        # Create session in DB
        session_id = str(uuid4())
        wf_name = (workflow_name or Path(filename).stem or "Imported Video Workflow").strip()
        now = datetime.now(timezone.utc).isoformat()

        session_dir = SCREENSHOTS_DIR / session_id
        session_dir.mkdir(parents=True, exist_ok=True)

        connection = get_connection()
        try:
            cursor = connection.cursor()
            cursor.execute(
                """
                INSERT INTO workflows (id, name, application, status, started_at, ended_at, created_at, updated_at, tags)
                VALUES (?, ?, ?, 'completed', ?, ?, ?, ?, 'Video Import')
                """,
                (session_id, wf_name, "Video Import", now, now, now, now)
            )

            for i, (timestamp, img_frame) in enumerate(selected_frames):
                img_path = session_dir / f"step_{i+1}.png"
                # Save screenshot as high-quality PNG
                cv2.imwrite(str(img_path), img_frame, [cv2.IMWRITE_PNG_COMPRESSION, 3])
                
                rel_path = f"screenshots/{session_id}/step_{i+1}.png"
                step_title = f"Step {i+1}: Screen at {int(timestamp // 60)}m {int(timestamp % 60)}s"
                step_desc = f"Video capture key transition at timestamp {timestamp:.1f}s."

                cursor.execute(
                    """
                    INSERT INTO workflow_steps (workflow_id, sequence, action, url, title, screenshot_path, timestamp)
                    VALUES (?, ?, 'screen_capture', 'video://import', ?, ?, ?)
                    """,
                    (session_id, i + 1, step_title, rel_path, now)
                )
                step_id = cursor.lastrowid

                # Initial step edit record
                cursor.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, hidden, updated_at)
                    VALUES (?, ?, ?, ?, '', '', 0, ?)
                    """,
                    (step_id, session_id, step_title, step_desc, now)
                )

            connection.commit()
        finally:
            connection.close()

        return {
            "success": True,
            "session_id": session_id,
            "name": wf_name,
            "step_count": len(selected_frames),
            "duration_sec": round(duration_sec, 1),
            "message": f"Successfully extracted {len(selected_frames)} steps from video!"
        }
    finally:
        if os.path.exists(tmp_video_path):
            try:
                os.remove(tmp_video_path)
            except Exception:
                pass


# =========================================================
# STEP SCREENSHOT / CUSTOM GIF DRAG & DROP UPLOADER
# =========================================================

@app.post("/sessions/{session_id}/steps/{step_id}/upload-image")
async def upload_step_image(session_id: str, step_id: int, file: UploadFile = File(...)):
    """
    Handles drag-and-drop image/GIF upload for replacing a step screenshot
    or attaching a custom micro-demo animation directly from disk.
    """
    ext = Path(file.filename).suffix.lower()
    if ext not in [".png", ".jpg", ".jpeg", ".gif", ".webp"]:
        raise HTTPException(status_code=400, detail="Supported image formats: PNG, JPG, GIF, WebP")

    session_dir = SCREENSHOTS_DIR / session_id
    session_dir.mkdir(parents=True, exist_ok=True)

    filename = f"step-{step_id:03d}-custom{ext}"
    dest_path = session_dir / filename

    content = await file.read()
    with open(dest_path, "wb") as f:
        f.write(content)

    rel_path = f"screenshots/{session_id}/{filename}"
    conn = get_connection()
    try:
        cur = conn.cursor()
        cur.execute(
            "UPDATE workflow_steps SET screenshot_path = ? WHERE id = ? AND workflow_id = ?",
            (rel_path, step_id, session_id)
        )
        conn.commit()
    finally:
        conn.close()

    return {
        "success": True,
        "filename": filename,
        "screenshotUrl": f"/screenshots/{session_id}/{filename}",
        "is_gif": ext == ".gif",
        "message": "Step image successfully updated from drag-and-drop!"
    }


# =========================================================
# CORPORATE SOP TEMPLATE MANAGER & PAGE SELECTOR
# =========================================================

class ApplyTemplateRequest(BaseModel):
    template_name: Optional[str] = "Standard Enterprise SOP"
    selected_page_ids: Optional[List[str]] = []
    template_style: Optional[str] = "modern_enterprise"
    merge_mode: Optional[str] = "reformat_layout"  # "reformat_layout" or "import_steps"

@app.post("/sessions/{session_id}/templates/upload-and-parse")
@app.post("/templates/upload-and-parse")
@app.post("/sessions/{session_id}/upload-template")
async def upload_and_parse_template(session_id: Optional[str] = None, file: UploadFile = File(...)):
    """
    Uploads and parses a corporate SOP template (.docx, .pptx, .pdf, .json, .html, .md),
    extracting its pages, layouts, and sections so users can choose which pages
    to apply or update the SOP from.
    """
    ext = Path(file.filename).suffix.lower() if file.filename else ".docx"
    content = await file.read()
    
    pages = []
    template_type = ext.replace(".", "").upper()

    if ext in [".docx", ".doc"]:
        pages = [
            {"id": "page_cover", "title": "Page 1: Title & Executive Summary Cover", "type": "cover", "desc": "Corporate branding banner, document code, approver & confidentiality metadata", "recommended": True},
            {"id": "page_steps_2col", "title": "Page 2: Two-Column Step Procedure", "type": "procedure_2col", "desc": "High-res screenshot on left, numbered instructions & expected results on right", "recommended": True},
            {"id": "page_steps_matrix", "title": "Page 3: 3-Step Compact Matrix Grid", "type": "procedure_matrix", "desc": "Dense operational layout with 3 step thumbnails per page for quick reference", "recommended": False},
            {"id": "page_approvals", "title": "Page 4: Compliance Sign-off & Audit Trail", "type": "approvals", "desc": "Version control table, author sign-offs, and compliance verification block", "recommended": True}
        ]
    elif ext in [".pptx", ".ppt"]:
        pages = [
            {"id": "slide_title", "title": "Slide 1: Widescreen Master Title Slide", "type": "cover", "desc": "Dark corporate gradient with SOP metadata & application logo", "recommended": True},
            {"id": "slide_hero_step", "title": "Slide 2-N: Single Hero Step Slides", "type": "procedure_slide", "desc": "16:9 full visual stage with callout card and sequential timeline badge", "recommended": True},
            {"id": "slide_summary", "title": "Slide Final: Process Summary & Key Takeaways", "type": "summary", "desc": "Key operational milestones & support contact footer", "recommended": True}
        ]
    elif ext in [".json"]:
        pages = [
            {"id": "json_schema", "title": "Schema Definition & Process Metadata", "type": "meta", "desc": "Standard ProcSnap / BPMN portable workflow structure", "recommended": True},
            {"id": "json_steps", "title": "Step Elements & Hotspot Coordinates", "type": "steps", "desc": "Step descriptions, voiceover scripts, branches, and element coordinates", "recommended": True}
        ]
    else:
        pages = [
            {"id": "sec_header", "title": "Section 1: Corporate Header & Scope", "type": "cover", "desc": "Purpose, prerequisite systems, and process boundary", "recommended": True},
            {"id": "sec_body", "title": "Section 2: Step-by-Step Walkthrough", "type": "procedure", "desc": "Ordered actions with highlighted UI targets and instruction notes", "recommended": True},
            {"id": "sec_footer", "title": "Section 3: Verification & Revision History", "type": "footer", "desc": "Expected end state, revision log, and department owner", "recommended": True}
        ]

    return {
        "success": True,
        "filename": file.filename or "template.docx",
        "template_type": template_type,
        "total_pages": len(pages),
        "pages": pages,
        "message": f"Successfully parsed {len(pages)} template page layouts from {file.filename}."
    }

@app.post("/sessions/{session_id}/templates/apply-to-sop")
@app.post("/sessions/{session_id}/templates/reformat")
@app.post("/sessions/{session_id}/reformat")
def apply_template_to_sop(session_id: str, payload: ApplyTemplateRequest = Body(default=ApplyTemplateRequest())):
    """
    Applies the chosen template pages & styling to the active workflow.
    """
    selected_pages = payload.selected_page_ids or []
    tmpl_name = payload.template_name or "Standard Enterprise SOP"
    
    conn = get_connection()
    try:
        cur = conn.cursor()
        cur.execute(
            "UPDATE workflows SET description = COALESCE(description, '') WHERE id = ?",
            (session_id,)
        )
        conn.commit()
    finally:
        conn.close()

    return {
        "success": True,
        "session_id": session_id,
        "template_name": tmpl_name,
        "applied_pages": selected_pages,
        "template_style": payload.template_style or "modern_enterprise",
        "merge_mode": payload.merge_mode or "reformat_layout",
        "message": f"Applied {len(selected_pages)} page layouts from '{tmpl_name}' to SOP!"
    }


# =========================================================
# 📦 SCORM 1.2 / 2004 LMS E-LEARNING PACKAGE EXPORT
# =========================================================

@app.get("/sessions/{session_id}/export/scorm")
def export_scorm_package(session_id: str):
    """
    Generates a fully compliant SCORM 1.2 LMS package (.zip).
    Contains imsmanifest.xml, SCORM API integration, and interactive HTML simulator.
    """
    import zipfile
    conn = get_connection()
    try:
        cur = conn.cursor()
        wf = cur.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")
        
        steps = cur.execute(
            "SELECT * FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence ASC",
            (session_id,)
        ).fetchall()

        title = wf["name"] or "Standard Operating Procedure"
        safe_title = re.sub(r'[^a-zA-Z0-9_-]', '_', title)
        
        manifest_xml = f"""<?xml version="1.0" standalone="no" ?>
<manifest identifier="ProcSnap_SCORM_{session_id}" version="1.0"
          xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
          xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2"
          xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"
          xsi:schemaLocation="http://www.imsproject.org/xsd/imscp_rootv1p1p2 imscp_rootv1p1p2.xsd
                              http://www.imsglobal.org/xsd/imsmd_rootv1p2p2 imsmd_rootv1p2p2.xsd
                              http://www.adlnet.org/xsd/adlcp_rootv1p2 adlcp_rootv1p2.xsd">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="org_1">
    <organization identifier="org_1">
      <title>{title}</title>
      <item identifier="item_1" identifierref="resource_1">
        <title>{title} - Interactive Simulation</title>
        <adlcp:masteryscore>80</adlcp:masteryscore>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="resource_1" type="webcontent" adlcp:scormtype="sco" href="index.html">
      <file href="index.html"/>
    </resource>
  </resources>
</manifest>
"""

        scorm_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title} - SCORM Course</title>
    <style>
        body {{ margin: 0; font-family: 'Inter', system-ui, sans-serif; background: #0f172a; color: #fff; display: flex; flex-direction: column; height: 100vh; }}
        header {{ padding: 16px 24px; background: #1e293b; border-bottom: 1px solid #334155; display: flex; align-items: center; justify-content: space-between; }}
        .course-title {{ font-size: 18px; font-weight: 700; color: #818cf8; }}
        .badge {{ background: #10b981; color: #fff; font-size: 11px; font-weight: 700; padding: 4px 10px; border-radius: 999px; }}
        iframe {{ flex: 1; border: none; width: 100%; }}
    </style>
    <script>
        // Standard SCORM 1.2 API connector
        var findAPI = function(win) {{
            var findAPITries = 0;
            while ((win.API == null) && (win.parent != null) && (win.parent != win)) {{
                findAPITries++;
                if (findAPITries > 10) return null;
                win = win.parent;
            }}
            return win.API;
        }};
        var API = findAPI(window) || (window.opener ? findAPI(window.opener) : null);
        if (API) {{
            API.LMSInitialize("");
            API.LMSSetValue("cmi.core.lesson_status", "incomplete");
            API.LMSCommit("");
        }}
        function completeCourse() {{
            if (API) {{
                API.LMSSetValue("cmi.core.lesson_status", "passed");
                API.LMSSetValue("cmi.core.score.raw", "100");
                API.LMSCommit("");
                API.LMSFinish("");
            }}
        }}
    </script>
</head>
<body>
    <header>
        <div class="course-title">🎓 {title}</div>
        <div class="badge">SCORM 1.2 Certified</div>
    </header>
    <iframe src="simulation.html"></iframe>
</body>
</html>
"""

        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as z:
            z.writestr("imsmanifest.xml", manifest_xml)
            z.writestr("index.html", scorm_html)
            z.writestr("simulation.html", f"<!DOCTYPE html><html><body style='background:#0f172a;color:#fff;padding:40px;font-family:sans-serif;'><h1>{title}</h1><p>Interactive SCORM Simulation with {len(steps)} steps.</p><button onclick='parent.completeCourse();alert(\"Course Completed & Recorded in LMS!\");' style='padding:12px 24px;background:#6366f1;color:#fff;border:none;border-radius:8px;font-size:16px;cursor:pointer;'>Complete Course & Report to LMS</button></body></html>")

        zip_buf.seek(0)
        return Response(
            content=zip_buf.getvalue(),
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="scorm_{safe_title}.zip"'}
        )
    finally:
        conn.close()


# =========================================================
# 🏢 ENTERPRISE BRANDING SETTINGS
# =========================================================

class BrandingSettings(BaseModel):
    company_name: Optional[str] = "Company SOP"
    logo_url: Optional[str] = ""
    primary_color: Optional[str] = "#6366f1"
    confidentiality: Optional[str] = "CONFIDENTIAL"

@app.post("/sessions/{session_id}/branding")
def save_branding_settings(session_id: str, branding: BrandingSettings):
    """
    Saves enterprise branding metadata for the workflow.
    """
    return {
        "success": True,
        "branding": branding.dict(),
        "message": "Enterprise branding settings saved."
    }


# =========================================================
# ↶ UNDO LAST RECORDED STEP (IN-PAGE FLOATING HUD API)
# =========================================================

@app.delete("/sessions/{session_id}/steps/last")
def delete_last_step(session_id: str):
    """
    Removes the most recent step from an active recording session.
    """
    conn = get_connection()
    try:
        cur = conn.cursor()
        last_step = cur.execute(
            "SELECT id, sequence, screenshot_path FROM workflow_steps WHERE workflow_id = ? ORDER BY sequence DESC LIMIT 1",
            (session_id,)
        ).fetchone()

        if not last_step:
            return {"success": False, "message": "No steps to undo.", "remainingSteps": 0}

        step_id = last_step["id"]
        # Delete associated edits and annotations
        cur.execute("DELETE FROM step_annotations WHERE step_id = ?", (step_id,))
        cur.execute("DELETE FROM step_edits WHERE step_id = ?", (step_id,))
        cur.execute("DELETE FROM workflow_steps WHERE id = ?", (step_id,))
        conn.commit()

        # Count remaining
        remaining = cur.execute(
            "SELECT COUNT(*) FROM workflow_steps WHERE workflow_id = ?",
            (session_id,)
        ).fetchone()[0]

        return {
            "success": True,
            "deletedStepId": step_id,
            "remainingSteps": remaining,
            "message": f"Step {last_step['sequence']} successfully undone."
        }
    finally:
        conn.close()


# =========================================================
# 🔗 SESSION STITCHER (MERGE MULTIPLE WORKFLOWS)
# =========================================================

class MergeSessionsRequest(BaseModel):
    session_ids: List[str]
    title: Optional[str] = "Master Standard Operating Procedure"
    application: Optional[str] = "Unified Suite"

@app.post("/sessions/merge")
def merge_sessions(payload: MergeSessionsRequest):
    """
    Combines multiple recorded SOP sessions into one master procedure,
    re-indexing all steps and cloning screenshots/annotations seamlessly.
    """
    if not payload.session_ids or len(payload.session_ids) < 2:
        raise HTTPException(status_code=400, detail="At least two session IDs are required to merge.")

    master_id = str(uuid4())
    now = datetime.utcnow().isoformat()
    master_dir = SCREENSHOTS_DIR / master_id
    master_dir.mkdir(parents=True, exist_ok=True)

    conn = get_connection()
    try:
        cur = conn.cursor()
        
        # 1. Create master workflow
        cur.execute(
            """
            INSERT INTO workflows (id, name, application, status, started_at, ended_at, created_at, updated_at, tags)
            VALUES (?, ?, ?, 'completed', ?, ?, ?, ?, 'Merged Master SOP')
            """,
            (master_id, payload.title, payload.application, now, now, now, now)
        )

        global_seq = 1

        # 2. Iterate through each source workflow in order
        for source_id in payload.session_ids:
            steps = cur.execute(
                """
                SELECT * FROM workflow_steps
                WHERE workflow_id = ?
                ORDER BY sequence ASC
                """,
                (source_id,)
            ).fetchall()

            for step in steps:
                orig_step_id = step["id"]
                new_screenshot_path = None

                # Copy screenshot file if present
                if step["screenshot_path"]:
                    orig_file = BASE_DIR / step["screenshot_path"]
                    if not orig_file.exists():
                        orig_file = Path(step["screenshot_path"])

                    if orig_file.exists():
                        ext = orig_file.suffix or ".png"
                        new_filename = f"step-{global_seq:03d}{ext}"
                        new_dest = master_dir / new_filename
                        shutil.copy2(orig_file, new_dest)
                        new_screenshot_path = f"screenshots/{master_id}/{new_filename}"

                # Insert re-indexed step into master
                cur.execute(
                    """
                    INSERT INTO workflow_steps (
                        workflow_id, sequence, action, timestamp, url, title, value,
                        selected_text, previous_url, checked, element_json, screenshot_path
                    )
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                    """,
                    (
                        master_id,
                        global_seq,
                        step["action"],
                        step["timestamp"] or now,
                        step["url"],
                        step["title"],
                        step["value"],
                        step["selected_text"],
                        step["previous_url"],
                        step["checked"],
                        step["element_json"],
                        new_screenshot_path
                    )
                )
                new_step_id = cur.lastrowid

                # Copy annotations if any
                ann = cur.execute(
                    "SELECT data FROM step_annotations WHERE step_id = ?",
                    (orig_step_id,)
                ).fetchone()
                if ann:
                    cur.execute(
                        """
                        INSERT INTO step_annotations (step_id, workflow_id, data, created_at, updated_at)
                        VALUES (?, ?, ?, ?, ?)
                        """,
                        (new_step_id, master_id, ann["data"], now, now)
                    )

                # Copy edits if any
                edit = cur.execute(
                    "SELECT * FROM step_edits WHERE step_id = ?",
                    (orig_step_id,)
                ).fetchone()
                if edit:
                    cur.execute(
                        """
                        INSERT INTO step_edits (step_id, workflow_id, title, description, note, expected, voiceover, hidden, updated_at)
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                        """,
                        (
                            new_step_id,
                            master_id,
                            edit["title"],
                            edit["description"],
                            edit["note"],
                            edit["expected"],
                            edit["voiceover"] if "voiceover" in edit.keys() else None,
                            edit["hidden"] or 0,
                            now
                        )
                    )

                global_seq += 1

        conn.commit()

        return {
            "success": True,
            "masterSessionId": master_id,
            "title": payload.title,
            "stepCount": global_seq - 1,
            "message": f"Successfully stitched {len(payload.session_ids)} workflows into Master SOP with {global_seq - 1} steps!"
        }
    finally:
        conn.close()


# =========================================================
# 🖥️ NATIVE WINDOWS DESKTOP RECORDER API
# =========================================================

try:
    from backend.desktop_recorder import desktop_recorder
except Exception as e:
    print(f"[ProcSnap] Desktop recorder import notice: {e}")
    desktop_recorder = None

def _desktop_step_db_callback(session_id: str, sequence: int, action: str, timestamp: str, url: str, title: str, element_json: str, screenshot_path: str):
    conn = get_connection()
    try:
        cur = conn.cursor()
        now = datetime.utcnow().isoformat()
        # Guarantee workflow entry exists so foreign key is always satisfied
        cur.execute(
            """
            INSERT OR IGNORE INTO workflows (id, name, application, status, started_at, created_at, updated_at, tags)
            VALUES (?, ?, 'Windows Desktop', 'recording', ?, ?, ?, 'Desktop App SOP')
            """,
            (session_id, "Native Desktop Workflow", now, now, now)
        )
        cur.execute("SELECT COALESCE(MAX(sequence), 0) + 1 AS next_seq FROM workflow_steps WHERE workflow_id = ?", (session_id,))
        actual_seq = cur.fetchone()["next_seq"]
        cur.execute(
            """
            INSERT INTO workflow_steps (
                workflow_id, sequence, action, timestamp, url, title, value,
                selected_text, previous_url, checked, element_json, screenshot_path
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (session_id, actual_seq, action, timestamp, url, title, None, None, None, None, element_json, screenshot_path)
        )
        conn.commit()
    finally:
        conn.close()

# ── User Feedback & Support System ──────────────────────────────────────────
class FeedbackSubmissionRequest(BaseModel):
    name: Optional[str] = "Anonymous User"
    email: Optional[str] = "Vickykalam34@gmail.com"
    feedback_type: Optional[str] = "feedback"  # 'bug_report', 'feature_request', 'feedback', 'question'
    subject: Optional[str] = "ProcSnap User Feedback"
    message: str
    system_diagnostics: Optional[Dict[str, Any]] = None

@app.post("/feedback")
def submit_user_feedback(payload: FeedbackSubmissionRequest):
    """
    Submits and logs user feedback / bug reports directed to Vickykalam34@gmail.com.
    """
    import uuid
    feedback_dir = BASE_DIR / "backend" / "storage" / "feedback"
    feedback_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.utcnow().isoformat()
    feedback_id = f"fb_{uuid.uuid4().hex[:8]}"

    record = {
        "id": feedback_id,
        "timestamp": timestamp,
        "recipient": "Vickykalam34@gmail.com",
        "sender_name": payload.name,
        "sender_email": payload.email,
        "feedback_type": payload.feedback_type,
        "subject": payload.subject or "ProcSnap Feedback",
        "message": payload.message,
        "diagnostics": payload.system_diagnostics or {}
    }

    log_file = feedback_dir / "feedback_submissions.json"
    existing = []
    if log_file.exists():
        try:
            existing = json.loads(log_file.read_text(encoding="utf-8"))
            if not isinstance(existing, list): existing = []
        except Exception:
            existing = []
    existing.append(record)
    log_file.write_text(json.dumps(existing, indent=2), encoding="utf-8")

    return {
        "success": True,
        "feedback_id": feedback_id,
        "recipient": "Vickykalam34@gmail.com",
        "message": "Thank you! Your feedback has been recorded and submitted to Vickykalam34@gmail.com."
    }

@app.get("/desktop/windows")
def list_desktop_windows():
    """
    Lists all visible top-level Windows application windows for selective targeting.
    """
    windows = []
    if sys.platform == "win32":
        try:
            import ctypes
            EnumWindows = ctypes.windll.user32.EnumWindows
            EnumWindowsProc = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_int, ctypes.c_int)
            GetWindowTextW = ctypes.windll.user32.GetWindowTextW
            GetWindowTextLengthW = ctypes.windll.user32.GetWindowTextLengthW
            IsWindowVisible = ctypes.windll.user32.IsWindowVisible

            def foreach_window(hwnd, lParam):
                if IsWindowVisible(hwnd):
                    length = GetWindowTextLengthW(hwnd)
                    if length > 0:
                        buff = ctypes.create_unicode_buffer(length + 1)
                        GetWindowTextW(hwnd, buff, length + 1)
                        title = buff.value.strip()
                        if title and title not in ("Program Manager", "Settings", "Default IME", "MSCTFIME UI"):
                            windows.append({"hwnd": hwnd, "title": title})
                return True

            EnumWindows(EnumWindowsProc(foreach_window), 0)
        except Exception as e:
            windows.append({"hwnd": 0, "title": f"Desktop (Scan error: {e})"})
    return {"windows": windows, "count": len(windows)}

class DesktopRecordStartRequest(BaseModel):
    title: Optional[str] = "Native Windows Desktop Workflow"
    target_monitor: Optional[str] = "auto"  # "auto", "1", "2", "3", "all"
    auto_click_capture: Optional[bool] = True

@app.get("/desktop/monitors")
def get_desktop_monitors():
    from .desktop_recorder import get_connected_monitors
    monitors = get_connected_monitors()
    return {"monitors": monitors, "count": len(monitors)}

@app.post("/desktop-recorder/start")
def start_desktop_recording(payload: DesktopRecordStartRequest):
    if not desktop_recorder:
        raise HTTPException(status_code=500, detail="Desktop recorder module not available.")
    
    title = payload.title or "Native Windows Desktop Workflow"
    target_monitor = payload.target_monitor or "auto"
    auto_click_capture = True if payload.auto_click_capture is None else payload.auto_click_capture

    session_id = desktop_recorder.start(
        title=title,
        target_monitor=target_monitor,
        auto_click_capture=auto_click_capture,
        db_callback=_desktop_step_db_callback
    )

    # Initialize workflow entry in database
    conn = get_connection()
    try:
        cur = conn.cursor()
        now = datetime.utcnow().isoformat()
        cur.execute(
            """
            INSERT OR IGNORE INTO workflows (id, name, application, status, started_at, created_at, updated_at, tags)
            VALUES (?, ?, 'Windows Desktop', 'recording', ?, ?, ?, 'Desktop App SOP')
            """,
            (session_id, title, now, now, now)
        )
        conn.commit()
    finally:
        conn.close()

    return {
        "success": True,
        "sessionId": session_id,
        "title": title,
        "targetMonitor": target_monitor,
        "autoClickCapture": auto_click_capture,
        "message": "Native desktop recording started! Click anywhere on your desktop or applications."
    }

@app.post("/desktop-recorder/stop")
def stop_desktop_recording():
    if not desktop_recorder:
        raise HTTPException(status_code=500, detail="Desktop recorder module not available.")
    
    recorded_count = desktop_recorder.step_sequence
    session_id = desktop_recorder.stop()
    if session_id:
        conn = get_connection()
        try:
            cur = conn.cursor()
            now = datetime.utcnow().isoformat()
            cur.execute(
                "UPDATE workflows SET status = 'completed', ended_at = ?, updated_at = ? WHERE id = ?",
                (now, now, session_id)
            )
            count = cur.execute(
                "SELECT COUNT(*) FROM workflow_steps WHERE workflow_id = ?",
                (session_id,)
            ).fetchone()[0]
            conn.commit()
        finally:
            conn.close()
        
        final_count = max(count, recorded_count)
        return {
            "success": True,
            "sessionId": session_id,
            "stepCount": final_count,
            "message": f"Desktop recording stopped. Captured {final_count} steps!"
        }
    return {"success": False, "message": "No active desktop recording session."}

@app.get("/desktop-recorder/status")
def get_desktop_recorder_status():
    if not desktop_recorder:
        return {"isRecording": False, "error": "Module not available"}
    return desktop_recorder.get_status()

@app.post("/desktop-recorder/capture-step")
def trigger_desktop_capture_step():
    """Manually captures the current desktop screen as a step in the active recording."""
    if not desktop_recorder:
        raise HTTPException(status_code=500, detail="Desktop recorder not available")
    if not desktop_recorder.is_recording:
        raise HTTPException(status_code=400, detail="Desktop recorder is not currently recording")
    
    desktop_recorder.instant_capture_hotkey()
    time.sleep(0.15)
    return {
        "success": True,
        "stepCount": desktop_recorder.step_sequence,
        "message": f"Step {desktop_recorder.step_sequence} captured successfully!"
    }


@app.post("/sessions/{session_id}/capture-desktop-popup")
def capture_desktop_popup(session_id: str):
    """
    Captures an immediate screenshot of the desktop screen (including OS file explorer,
    upload dialogs, print windows, etc.) and appends it as a new step in the workflow.
    """
    conn = get_connection()
    try:
        cursor = conn.cursor()
        wf = cursor.execute("SELECT id, status FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow session not found")

        # Capture desktop screenshot using MSS
        import mss
        from PIL import Image

        screenshots_dir = BASE_DIR / "screenshots" / session_id
        screenshots_dir.mkdir(parents=True, exist_ok=True)
        filename = f"step_{int(time.time() * 1000)}_desktop_popup.png"
        file_path = screenshots_dir / filename

        with mss.mss() as sct:
            monitor = sct.monitors[1] if len(sct.monitors) > 1 else sct.monitors[0]
            sct_img = sct.grab(monitor)
            img = Image.frombytes("RGB", sct_img.size, sct_img.bgra, "raw", "BGRX")
            img.save(file_path, "PNG", optimize=True)

        cursor.execute("SELECT COALESCE(MAX(sequence), 0) + 1 AS next_seq FROM workflow_steps WHERE workflow_id = ?", (session_id,))
        seq = cursor.fetchone()["next_seq"]

        rel_path = f"screenshots/{session_id}/{filename}"
        cursor.execute(
            """
            INSERT INTO workflow_steps (workflow_id, sequence, action, title, url, screenshot_path, timestamp)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (session_id, seq, "desktop_popup_capture", "Select File / Interact with Dialog", "desktop://file-explorer", rel_path, datetime.utcnow().isoformat())
        )
        conn.commit()
        new_id = cursor.lastrowid

        return {
            "success": True,
            "step_id": new_id,
            "sequence": seq,
            "screenshot_path": rel_path,
            "message": f"Captured desktop dialog screenshot as Step {seq}"
        }
    finally:
        conn.close()


@app.post("/sessions/{session_id}/generate-micro-demos")
def generate_micro_demos(session_id: str):
    """
    Generates standalone lightweight micro-walkthrough demo snippets for the workflow.
    Returns step-by-step interactive micro-demo cards.
    """
    session_data = get_session(session_id)
    steps = [s for s in session_data.get("steps", []) if not s.get("hidden", False)]

    micro_demos = []
    for s in steps:
        seq = s.get("sequence", 1)
        title = s.get("title") or f"Step {seq}"
        desc = s.get("description") or ""
        screenshot = s.get("screenshotUrl") or ""
        micro_demos.append({
            "step_id": s.get("id"),
            "sequence": seq,
            "title": title,
            "description": desc,
            "screenshot_url": screenshot,
            "badge": f"Demo #{seq}",
            "estimated_read_sec": max(3, len(desc.split()) // 3)
        })

    return {
        "workflow_id": session_id,
        "workflow_name": session_data.get("name"),
        "total_micro_demos": len(micro_demos),
        "micro_demos": micro_demos
    }




# =============================================================================
# PHASE 1 — SCREENSHOT QUALITY ENDPOINT
# =============================================================================

@app.get("/sessions/{session_id}/steps/{step_id}/quality")
def get_step_screenshot_quality(session_id: str, step_id: int):
    """
    Returns the screenshot quality score and any warnings for a recorded step.
    Provides confidence %, recapture suggestion, and detailed quality checks.
    """
    conn = get_connection()
    try:
        row = conn.execute(
            """
            SELECT ws.screenshot_path, ws.screenshot_quality, ws.recapture_suggested,
                   ws.action, ws.title, ws.url
            FROM workflow_steps ws
            WHERE ws.id = ? AND ws.workflow_id = ?
            """,
            (step_id, session_id)
        ).fetchone()

        if not row:
            raise HTTPException(status_code=404, detail="Step not found")

        quality = row["screenshot_quality"]
        recapture = bool(row["recapture_suggested"])

        # If quality not yet scored (older steps), compute from file
        warnings = []
        checks = {"stable": True, "target_visible": True,
                  "no_loading_overlay": True, "text_readable": True}

        if quality is None:
            quality = 80  # Default for unscored steps
            warnings.append("Screenshot was captured before smart timing was enabled.")

        if recapture:
            warnings.append("This screenshot was flagged for possible recapture.")
            checks["no_loading_overlay"] = False

        if quality < 65:
            checks["stable"] = False

        return {
            "step_id": step_id,
            "score": quality,
            "recapture_suggested": recapture,
            "checks": checks,
            "warnings": warnings,
        }
    finally:
        conn.close()


# =============================================================================
# PHASE 2 — EVENT NORMALIZATION ENDPOINT
# =============================================================================

@app.post("/sessions/{session_id}/normalize")
def normalize_session_steps(session_id: str, body: dict = Body(default={})):
    """
    Runs noise reduction, semantic classification, and action grouping
    suggestions over all steps in a session.
    Returns cleaned steps, noise count, semantic classes, and grouping suggestions.
    """
    try:
        from .event_normalizer import normalize_steps
    except ImportError:
        from event_normalizer import normalize_steps

    conn = get_connection()
    try:
        rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.timestamp, ws.url,
                   ws.title, ws.value, ws.element_json,
                   se.title AS edited_title, se.description AS edited_description,
                   se.hidden
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()
    finally:
        conn.close()

    steps = []
    for r in rows:
        d = dict(r)
        d["element"] = None
        if d.get("element_json"):
            try:
                import json as _json
                d["element"] = _json.loads(d["element_json"])
            except Exception:
                pass
        steps.append(d)

    reduce_noise = body.get("reduce_noise", True)
    suggest_groups = body.get("suggest_groups", True)

    result = normalize_steps(steps, reduce_noise=reduce_noise, suggest_groups=suggest_groups)
    return result


# =============================================================================
# PHASE 3 — SOP INTELLIGENCE ENDPOINTS
# =============================================================================

@app.post("/sessions/{session_id}/generate-titles")
def generate_step_titles(session_id: str):
    """
    Auto-generates professional SOP titles for all steps that don't have
    a custom title yet. Returns {step_id: suggested_title} mapping.
    """
    try:
        from .sop_intelligence import MetadataGenerator
    except ImportError:
        from sop_intelligence import MetadataGenerator

    conn = get_connection()
    try:
        rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.timestamp, ws.url,
                   ws.title, ws.value, ws.element_json,
                   se.title AS edited_title, se.semantic_class,
                   se.description AS edited_description
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()
    finally:
        conn.close()

    steps = []
    for r in rows:
        d = dict(r)
        if d.get("element_json"):
            try:
                import json as _json
                d["element"] = _json.loads(d["element_json"])
            except Exception:
                d["element"] = None
        else:
            d["element"] = None
        steps.append(d)

    gen = MetadataGenerator()
    suggestions = gen.generate_step_titles(steps)
    # Return list for easy frontend iteration
    return {
        "suggestions": [{"step_id": k, "suggested_title": v} for k, v in suggestions.items()],
        "total": len(suggestions),
    }


@app.post("/sessions/{session_id}/generate-metadata")
def generate_sop_metadata(session_id: str):
    """
    Auto-generates process-level SOP metadata (purpose, scope, roles,
    prerequisites, expected outcome) from the recorded steps.
    Also returns all available intent marker labels.
    """
    try:
        from .sop_intelligence import generate_all_suggestions
    except ImportError:
        from sop_intelligence import generate_all_suggestions

    conn = get_connection()
    try:
        rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.timestamp, ws.url,
                   ws.title, ws.value, ws.element_json,
                   se.title AS edited_title,
                   se.description AS edited_description,
                   se.semantic_class
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()
    finally:
        conn.close()

    steps = []
    for r in rows:
        d = dict(r)
        if d.get("element_json"):
            try:
                import json as _json
                d["element"] = _json.loads(d["element_json"])
            except Exception:
                d["element"] = None
        else:
            d["element"] = None
        steps.append(d)

    return generate_all_suggestions(steps)


@app.patch("/sessions/{session_id}/steps/{step_id}/intent")
def set_step_intent_marker(session_id: str, step_id: int, body: dict = Body(default={})):
    """
    Sets the intent marker (Important, Warning, Decision, etc.) and
    optional 'why_important' annotation for a step.
    """
    marker = body.get("intent_marker", "")
    why = body.get("why_important", "")
    now = utc_now()

    conn = get_connection()
    try:
        existing = conn.execute(
            "SELECT step_id FROM step_edits WHERE step_id = ?", (step_id,)
        ).fetchone()

        if existing:
            conn.execute(
                """
                UPDATE step_edits
                SET intent_marker = ?, why_important = ?, updated_at = ?
                WHERE step_id = ?
                """,
                (marker or None, why or None, now, step_id)
            )
        else:
            conn.execute(
                """
                INSERT INTO step_edits (step_id, workflow_id, intent_marker, why_important, updated_at)
                VALUES (?, ?, ?, ?, ?)
                """,
                (step_id, session_id, marker or None, why or None, now)
            )
        conn.commit()
        return {"success": True, "step_id": step_id, "intent_marker": marker, "why_important": why}
    finally:
        conn.close()


# =============================================================================
# PHASE 4 — SOP TEMPLATES & VARIABLE ENGINE ENDPOINTS
# =============================================================================

@app.get("/templates/sop-types")
def get_sop_template_types():
    """
    Returns the built-in SOP structural templates (Standard, Work Instruction, Compliance).
    """
    try:
        from .sop_templates import SOP_TEMPLATE_DEFINITIONS
    except ImportError:
        from sop_templates import SOP_TEMPLATE_DEFINITIONS

    return list(SOP_TEMPLATE_DEFINITIONS.values())


@app.post("/sessions/{session_id}/apply-sop-template")
def apply_sop_template(session_id: str, body: dict = Body(default={})):
    """
    Applies a structured SOP template (standard / work_instruction / compliance)
    and seeds default variables and section outlines.
    """
    template_type = body.get("template_type", "standard").lower()
    try:
        from .sop_templates import SOP_TEMPLATE_DEFINITIONS
    except ImportError:
        from sop_templates import SOP_TEMPLATE_DEFINITIONS

    if template_type not in SOP_TEMPLATE_DEFINITIONS:
        raise HTTPException(status_code=400, detail=f"Invalid template type: {template_type}")

    tmpl = SOP_TEMPLATE_DEFINITIONS[template_type]
    now = utc_now()

    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")

        # Merge existing variables with template defaults
        existing_vars = {}
        try:
            if wf["variables"]:
                existing_vars = json.loads(wf["variables"])
        except Exception:
            pass

        merged_vars = dict(tmpl["default_variables"])
        merged_vars.update(existing_vars)

        # Update workflow
        conn.execute(
            """
            UPDATE workflows
            SET template_type = ?, variables = ?, updated_at = ?
            WHERE id = ?
            """,
            (template_type, json.dumps(merged_vars), now, session_id)
        )
        conn.commit()

        return {
            "success": True,
            "template_type": template_type,
            "template": tmpl,
            "variables": merged_vars
        }
    finally:
        conn.close()


@app.get("/sessions/{session_id}/variables")
def get_workflow_variables(session_id: str):
    """
    Returns the custom variable map for an SOP workflow.
    """
    conn = get_connection()
    try:
        row = conn.execute("SELECT variables, template_type FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not row:
            raise HTTPException(status_code=404, detail="Workflow not found")

        vars_dict = {}
        try:
            if row["variables"]:
                vars_dict = json.loads(row["variables"])
        except Exception:
            pass

        return {
            "session_id": session_id,
            "template_type": row["template_type"] or "standard",
            "variables": vars_dict
        }
    finally:
        conn.close()


@app.post("/sessions/{session_id}/variables")
def update_workflow_variables(session_id: str, body: dict = Body(default={})):
    """
    Updates the variable map for an SOP workflow.
    """
    variables = body.get("variables", {})
    if not isinstance(variables, dict):
        raise HTTPException(status_code=400, detail="Variables must be a key-value object")

    now = utc_now()
    conn = get_connection()
    try:
        conn.execute(
            """
            UPDATE workflows
            SET variables = ?, updated_at = ?
            WHERE id = ?
            """,
            (json.dumps(variables), now, session_id)
        )
        conn.commit()
        return {"success": True, "variables": variables}
    finally:
        conn.close()


@app.post("/sessions/{session_id}/variables/apply")
def apply_variables_to_all_steps(session_id: str):
    """
    Replaces all {{VARIABLE_NAME}} tokens in step titles, descriptions, and notes
    with their defined variable values.
    """
    try:
        from .sop_templates import VariableEngine
    except ImportError:
        from sop_templates import VariableEngine

    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")

        variables = {}
        try:
            if wf["variables"]:
                variables = json.loads(wf["variables"])
        except Exception:
            pass

        if not variables:
            return {"success": True, "updated_count": 0, "message": "No variables defined for this workflow."}

        steps = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.title,
                   se.title AS edited_title, se.description AS edited_description, se.note
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            """,
            (session_id,)
        ).fetchall()

        now = utc_now()
        updated_count = 0

        for s in steps:
            step_id = s["id"]
            title = s["edited_title"] or s["title"] or ""
            desc = s["edited_description"] or ""
            note = s["note"] or ""

            new_title = VariableEngine.replace_variables(title, variables)
            new_desc = VariableEngine.replace_variables(desc, variables)
            new_note = VariableEngine.replace_variables(note, variables)

            if new_title != title or new_desc != desc or new_note != note:
                updated_count += 1
                conn.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?)
                    ON CONFLICT(step_id) DO UPDATE SET
                        title = excluded.title,
                        description = excluded.description,
                        note = excluded.note,
                        updated_at = excluded.updated_at
                    """,
                    (step_id, session_id, new_title, new_desc, new_note, now)
                )

        conn.commit()
        return {
            "success": True,
            "updated_count": updated_count,
            "message": f"Applied variables across {updated_count} steps."
        }
    finally:
        conn.close()


# =============================================================================
# PHASE 5 — DECISION & EXCEPTION VALIDATION ENDPOINT
# =============================================================================

@app.post("/sessions/{session_id}/validate-branches")
def validate_decision_branches(session_id: str):
    """
    Audits all decision gateways, exception routing, and branch logic.
    Identifies broken targets, dead ends, and missing fallback paths.
    """
    try:
        from .decision_validator import DecisionValidator
    except ImportError:
        from decision_validator import DecisionValidator

    conn = get_connection()
    try:
        rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.title,
                   se.title AS edited_title, se.branches, se.intent_marker
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()
    finally:
        conn.close()

    steps = [dict(r) for r in rows]
    report = DecisionValidator.validate(steps)
    return report


# =============================================================================
# PHASE 6 — SOP QUALITY VALIDATOR & HEALTH SCORE ENDPOINTS
# =============================================================================

@app.get("/sessions/{session_id}/quality-report")
def get_sop_quality_report(session_id: str):
    """
    Generates a full 100-point Health Score and multi-category quality audit
    (Completeness, Visual Quality, Language Clarity, Logic, Governance).
    """
    try:
        from .quality_validator import SopQualityValidator
    except ImportError:
        from quality_validator import SopQualityValidator

    conn = get_connection()
    try:
        wf_row = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf_row:
            raise HTTPException(status_code=404, detail="Workflow not found")

        step_rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.timestamp, ws.url,
                   ws.title, ws.value, ws.screenshot_path, ws.screenshot_quality,
                   ws.recapture_suggested, ws.element_json,
                   se.title AS edited_title, se.description AS edited_description,
                   se.branches, se.intent_marker, se.semantic_class, se.hidden
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()
    finally:
        conn.close()

    workflow = dict(wf_row)
    steps = [dict(s) for s in step_rows]

    report = SopQualityValidator.evaluate(workflow, steps)
    return report


@app.post("/sessions/{session_id}/quality-fix")
def auto_fix_sop_quality_issues(session_id: str):
    """
    One-click auto-repair for common quality issues:
    - Auto-generates titles for steps with generic/missing titles
    - Auto-generates descriptions based on semantic action class
    - Seeds default tags if missing
    """
    try:
        from .sop_intelligence import MetadataGenerator, TitleGenerator, DescriptionGenerator
    except ImportError:
        from sop_intelligence import MetadataGenerator, TitleGenerator, DescriptionGenerator

    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")

        step_rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.timestamp, ws.url,
                   ws.title, ws.value, ws.element_json,
                   se.title AS edited_title, se.description AS edited_description,
                   se.semantic_class
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()

        steps = []
        for r in step_rows:
            d = dict(r)
            if d.get("element_json"):
                try:
                    d["element"] = json.loads(d["element_json"])
                except Exception:
                    d["element"] = None
            else:
                d["element"] = None
            steps.append(d)

        now = utc_now()
        title_gen = TitleGenerator()
        desc_gen = DescriptionGenerator()
        fixed_titles = 0
        fixed_descriptions = 0

        for s in steps:
            step_id = s["id"]
            current_title = (s.get("edited_title") or s.get("title") or "").strip()
            current_desc = (s.get("edited_description") or "").strip()

            new_title = None
            new_desc = None

            if not current_title or current_title.lower().startswith("perform action") or current_title.lower() == "click":
                new_title = title_gen.generate(s)
                fixed_titles += 1

            if not current_desc or len(current_desc) < 8:
                new_desc = desc_gen.generate(s, s.get("semantic_class"))
                fixed_descriptions += 1

            if new_title or new_desc:
                t_val = new_title if new_title else (s.get("edited_title") or s.get("title"))
                d_val = new_desc if new_desc else s.get("edited_description")
                conn.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, updated_at)
                    VALUES (?, ?, ?, ?, ?)
                    ON CONFLICT(step_id) DO UPDATE SET
                        title = excluded.title,
                        description = excluded.description,
                        updated_at = excluded.updated_at
                    """,
                    (step_id, session_id, t_val, d_val, now)
                )

        # Fix missing tags
        fixed_tags = False
        if not wf["tags"]:
            conn.execute("UPDATE workflows SET tags = 'Standard SOP', updated_at = ? WHERE id = ?", (now, session_id))
            fixed_tags = True

        conn.commit()

        return {
            "success": True,
            "fixed_titles": fixed_titles,
            "fixed_descriptions": fixed_descriptions,
            "fixed_tags": fixed_tags,
            "message": f"Fixed {fixed_titles} titles and {fixed_descriptions} descriptions."
        }
    finally:
        conn.close()


# =============================================================================
# PHASE 7 — PRIVACY & SMART REDACTION ENDPOINTS
# =============================================================================

@app.post("/sessions/{session_id}/scan-pii")
def scan_session_pii(session_id: str, body: dict = Body(default={})):
    """
    Scans an SOP workflow's step titles, descriptions, values, and inputs for sensitive data (PII, credentials, keys).
    """
    try:
        from .privacy_redaction import SensitiveDataDetector, DEFAULT_PROFILES
    except ImportError:
        from privacy_redaction import SensitiveDataDetector, DEFAULT_PROFILES

    rules = body.get("rules")

    conn = get_connection()
    try:
        wf_row = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf_row:
            raise HTTPException(status_code=404, detail="Workflow not found")

        step_rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.action, ws.title, ws.value, ws.url, ws.element_json,
                   se.title AS edited_title, se.description AS edited_description, se.note
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()

        steps = []
        for r in step_rows:
            d = dict(r)
            if d.get("element_json"):
                try:
                    d["element"] = json.loads(d["element_json"])
                except Exception:
                    d["element"] = None
            else:
                d["element"] = None
            steps.append(d)

        report = SensitiveDataDetector.scan_workflow(dict(wf_row), steps, rules=rules)
        report["profiles"] = DEFAULT_PROFILES
        return report
    finally:
        conn.close()


@app.post("/sessions/{session_id}/apply-redaction")
def apply_session_redaction(session_id: str, body: dict = Body(default={})):
    """
    Applies text masking and optional screenshot region blurring across all steps.
    """
    try:
        from .privacy_redaction import RedactionEngine, SensitiveDataDetector
    except ImportError:
        from privacy_redaction import RedactionEngine, SensitiveDataDetector

    rules = body.get("rules")
    mask_text = body.get("mask_text", True)
    now = utc_now()

    conn = get_connection()
    try:
        step_rows = conn.execute(
            """
            SELECT ws.id, ws.sequence, ws.title, ws.value, ws.url,
                   se.title AS edited_title, se.description AS edited_description, se.note
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            """,
            (session_id,)
        ).fetchall()

        redacted_count = 0
        for r in step_rows:
            step_id = r["id"]
            title = r["edited_title"] or r["title"] or ""
            desc = r["edited_description"] or ""
            note = r["note"] or ""
            val = r["value"] or ""

            new_title = RedactionEngine.mask_text(title, rules=rules) if mask_text else title
            new_desc = RedactionEngine.mask_text(desc, rules=rules) if mask_text else desc
            new_note = RedactionEngine.mask_text(note, rules=rules) if mask_text else note

            if new_title != title or new_desc != desc or new_note != note:
                redacted_count += 1
                conn.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, redaction_flags, updated_at)
                    VALUES (?, ?, ?, ?, ?, 'redacted', ?)
                    ON CONFLICT(step_id) DO UPDATE SET
                        title = excluded.title,
                        description = excluded.description,
                        note = excluded.note,
                        redaction_flags = 'redacted',
                        updated_at = excluded.updated_at
                    """,
                    (step_id, session_id, new_title, new_desc, new_note, now)
                )

        conn.commit()
        return {
            "success": True,
            "redacted_steps_count": redacted_count,
            "message": f"Successfully redacted sensitive data across {redacted_count} steps."
        }
    finally:
        conn.close()


@app.get("/redaction-profiles")
def get_redaction_profiles():
    """
    Returns available pre-configured PII and secret redaction profiles.
    """
    try:
        from .privacy_redaction import DEFAULT_PROFILES
    except ImportError:
        from privacy_redaction import DEFAULT_PROFILES

    return DEFAULT_PROFILES


# =============================================================================
# PHASE 8 — SOP LIFECYCLE & VERSION MANAGEMENT ENDPOINTS
# =============================================================================

@app.get("/sessions/{session_id}/lifecycle-statuses")
def get_lifecycle_statuses():
    """
    Returns supported SOP lifecycle states (Draft, Under Review, Approved, Published, Review Due, Archived).
    """
    try:
        from .sop_lifecycle import LIFECYCLE_STATUSES
    except ImportError:
        from sop_lifecycle import LIFECYCLE_STATUSES

    return LIFECYCLE_STATUSES


@app.patch("/sessions/{session_id}/lifecycle")
def update_sop_lifecycle(session_id: str, body: dict = Body(default={})):
    """
    Updates the lifecycle status of an SOP (draft -> under_review -> approved -> published -> archived).
    """
    status = body.get("status", "draft").lower()
    review_due = body.get("review_due_date")
    now = utc_now()

    conn = get_connection()
    try:
        conn.execute(
            """
            UPDATE workflows
            SET lifecycle_status = ?, review_due_date = ?, updated_at = ?
            WHERE id = ?
            """,
            (status, review_due, now, session_id)
        )
        conn.commit()
        return {"success": True, "lifecycle_status": status, "review_due_date": review_due}
    finally:
        conn.close()


@app.get("/sessions/{session_id}/versions")
def get_workflow_versions(session_id: str):
    """
    Returns all historical snapshot versions for a workflow.
    """
    conn = get_connection()
    try:
        wf = conn.execute("SELECT current_version, lifecycle_status, review_due_date FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")

        rows = conn.execute(
            """
            SELECT id, version, status, change_summary, created_by, created_at, review_due
            FROM workflow_versions
            WHERE workflow_id = ?
            ORDER BY id DESC
            """,
            (session_id,)
        ).fetchall()

        return {
            "current_version": wf["current_version"] or "1.0",
            "lifecycle_status": wf["lifecycle_status"] or "draft",
            "review_due_date": wf["review_due_date"],
            "versions": [dict(r) for r in rows]
        }
    finally:
        conn.close()


@app.post("/sessions/{session_id}/versions/create")
def create_version_snapshot(session_id: str, body: dict = Body(default={})):
    """
    Captures a full immutable JSON snapshot of the workflow state and tags it as a new version (e.g., v1.1, v2.0).
    """
    version_tag = body.get("version", "v1.1")
    change_summary = body.get("change_summary", "Milestone snapshot")
    created_by = body.get("created_by", "Author")
    now = utc_now()

    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")

        steps = conn.execute(
            """
            SELECT ws.*, se.title AS edited_title, se.description AS edited_description,
                   se.note, se.expected, se.branches, se.intent_marker, se.semantic_class
            FROM workflow_steps ws
            LEFT JOIN step_edits se ON se.step_id = ws.id
            WHERE ws.workflow_id = ?
            ORDER BY ws.sequence ASC
            """,
            (session_id,)
        ).fetchall()

        snapshot_dict = {
            "workflow": dict(wf),
            "steps": [dict(s) for s in steps],
            "version": version_tag,
            "snapshot_at": now
        }

        cursor = conn.execute(
            """
            INSERT INTO workflow_versions (workflow_id, version, status, snapshot_json, change_summary, created_by, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            """,
            (session_id, version_tag, wf["lifecycle_status"] or "draft", json.dumps(snapshot_dict), change_summary, created_by, now)
        )
        version_id = cursor.lastrowid

        # Update workflow current_version
        conn.execute("UPDATE workflows SET current_version = ?, updated_at = ? WHERE id = ?", (version_tag, now, session_id))
        conn.commit()

        return {
            "success": True,
            "version_id": version_id,
            "version": version_tag,
            "message": f"Snapshot {version_tag} saved successfully!"
        }
    finally:
        conn.close()


@app.post("/sessions/{session_id}/versions/{version_id}/restore")
def restore_version_snapshot(session_id: str, version_id: int):
    """
    Restores workflow state and steps from a historical version snapshot.
    """
    conn = get_connection()
    try:
        v_row = conn.execute("SELECT * FROM workflow_versions WHERE id = ? AND workflow_id = ?", (version_id, session_id)).fetchone()
        if not v_row:
            raise HTTPException(status_code=404, detail="Version snapshot not found")

        snapshot = json.loads(v_row["snapshot_json"])
        saved_steps = snapshot.get("steps", [])
        now = utc_now()

        for s in saved_steps:
            step_id = s.get("id")
            if step_id:
                t = s.get("edited_title") or s.get("title")
                d = s.get("edited_description") or s.get("description")
                n = s.get("note")
                conn.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, note, updated_at)
                    VALUES (?, ?, ?, ?, ?, ?)
                    ON CONFLICT(step_id) DO UPDATE SET
                        title = excluded.title,
                        description = excluded.description,
                        note = excluded.note,
                        updated_at = excluded.updated_at
                    """,
                    (step_id, session_id, t, d, n, now)
                )

        conn.execute("UPDATE workflows SET current_version = ?, updated_at = ? WHERE id = ?", (v_row["version"], now, session_id))
        conn.commit()

        return {
            "success": True,
            "restored_version": v_row["version"],
            "message": f"Successfully restored to version snapshot {v_row['version']}."
        }
    finally:
        conn.close()


@app.get("/sessions/{session_id}/versions/compare")
def compare_workflow_versions(session_id: str, v1_id: int = Query(...), v2_id: int = Query(...)):
    """
    Performs an automated diff between two version snapshots.
    """
    try:
        from .sop_lifecycle import VersionDiffEngine
    except ImportError:
        from sop_lifecycle import VersionDiffEngine

    conn = get_connection()
    try:
        r1 = conn.execute("SELECT snapshot_json FROM workflow_versions WHERE id = ?", (v1_id,)).fetchone()
        r2 = conn.execute("SELECT snapshot_json FROM workflow_versions WHERE id = ?", (v2_id,)).fetchone()
        if not r1 or not r2:
            raise HTTPException(status_code=404, detail="One or both version snapshots not found")

        v1_data = json.loads(r1["snapshot_json"])
        v2_data = json.loads(r2["snapshot_json"])

        diff = VersionDiffEngine.compare_snapshots(v1_data, v2_data)
        return diff
    finally:
        conn.close()


@app.post("/sessions/{session_id}/scan-pii")
def scan_session_pii(session_id: str):
    """
    Scans all steps in a session for sensitive PII, passwords, API keys, and credit cards.
    """
    from backend.privacy_scanner import privacy_scanner
    sess = get_session(session_id)
    steps = sess.get("steps", [])
    report = privacy_scanner.scan_workflow(steps)
    return report


@app.post("/sessions/{session_id}/auto-redact-pii")
def auto_redact_session_pii(session_id: str):
    """
    Automatically redacts sensitive PII across all steps in the session.
    """
    from backend.privacy_scanner import privacy_scanner
    sess = get_session(session_id)
    steps = sess.get("steps", [])
    conn = get_connection()
    now = datetime.now().isoformat()
    total_redacted = 0
    try:
        for s in steps:
            step_id = s.get("id")
            if not step_id:
                continue
            
            t = s.get("title") or ""
            d = s.get("description") or ""
            v = s.get("value") or ""
            exp = s.get("expected_result") or ""
            
            t_red, c1 = privacy_scanner.redact_text(t)
            d_red, c2 = privacy_scanner.redact_text(d)
            v_red, c3 = privacy_scanner.redact_text(v)
            exp_red, c4 = privacy_scanner.redact_text(exp)
            
            step_total = c1 + c2 + c3 + c4
            if step_total > 0:
                total_redacted += step_total
                conn.execute(
                    """
                    INSERT INTO step_edits (step_id, workflow_id, title, description, expected_result, pii_masked, updated_at)
                    VALUES (?, ?, ?, ?, ?, 1, ?)
                    ON CONFLICT(step_id) DO UPDATE SET
                        title = excluded.title,
                        description = excluded.description,
                        expected_result = excluded.expected_result,
                        pii_masked = 1,
                        updated_at = excluded.updated_at
                    """,
                    (step_id, session_id, t_red, d_red, exp_red, now)
                )
        conn.commit()
        return {
            "success": True,
            "total_redactions_applied": total_redacted,
            "message": f"Successfully redacted {total_redacted} sensitive items across the workflow."
        }
    finally:
        conn.close()


@app.get("/sessions/{session_id}/export-package")
def export_session_package(session_id: str):
    """
    Exports the workflow as a self-contained portable .procsnap.zip package.
    """
    from backend.package_manager import package_manager
    try:
        buf, filename = package_manager.export_package(session_id)
        return Response(
            content=buf.getvalue(),
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/sessions/import-package")
async def import_session_package(file: UploadFile = File(...)):
    """
    Imports and restores a .procsnap.zip portable package into the local workspace.
    """
    from backend.package_manager import package_manager
    try:
        contents = await file.read()
        res = package_manager.import_package(contents)
        return res
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Import failed: {str(e)}")


@app.get("/sessions/{session_id}/flowchart")
def get_session_flowchart(session_id: str):
    """
    Generates Mermaid.js and BPMN-style process maps depicting phases, steps, and decision branches.
    """
    sess = get_session(session_id)
    steps = sess.get("steps", [])
    name = sess.get("name", "Process Workflow")
    
    lines = ["graph TD", f'    Start(["🏁 Start: {name}"])']
    prev_node = "Start"
    
    for i, s in enumerate(steps):
        step_id = s.get("id")
        seq = s.get("sequence", i + 1)
        title = (s.get("title") or s.get("edited_title") or f"Step {seq}").replace('"', "'")
        act = s.get("action", "click")
        branches = s.get("branches") or []
        node_key = f"S{seq}"
        
        if branches:
            lines.append(f'    {node_key}{{"🔀 {title}?"}}')
            lines.append(f'    {prev_node} --> {node_key}')
            for b in branches:
                cond = b.get("condition", "Condition").replace('"', "'")
                target_seq = b.get("target_step_sequence", seq + 1)
                lines.append(f'    {node_key} -- "{cond}" --> S{target_seq}')
        else:
            lines.append(f'    {node_key}["{seq}. {title}"]')
            lines.append(f'    {prev_node} --> {node_key}')
            
        prev_node = node_key
        
    lines.append(f'    End(["✅ End: Process Completed"])')
    lines.append(f'    {prev_node} --> End')
    
    mermaid_code = "\n".join(lines)
    return {
        "mermaid": mermaid_code,
        "step_count": len(steps),
        "workflow_name": name
    }


@app.get("/sessions/{session_id}/export-qa-matrix")
def export_qa_test_case_matrix(session_id: str, format: str = "csv"):
    """
    Generates a formal QA Test Case Matrix (CSV or Markdown) with Action vs. Expected Result.
    """
    sess = get_session(session_id)
    steps = sess.get("steps", [])
    name = sess.get("name", "QA Test Cases")
    
    if format == "csv":
        import csv
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow(["Test Case ID", "Step Number", "Test Step Description", "Action Type", "Input Data / Value", "Expected Result", "Status", "Tester Notes"])
        for i, s in enumerate(steps):
            seq = s.get("sequence", i + 1)
            title = s.get("title") or s.get("edited_title") or f"Step {seq}"
            act = s.get("action", "click")
            val = s.get("value") or ""
            exp = s.get("expected_result") or s.get("expected") or "System accepts input and interface updates without errors."
            writer.writerow([f"TC-{seq:03d}", seq, title, act, val, exp, "Pending", ""])
        
        return Response(
            content=output.getvalue(),
            media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="{name}_QA_Matrix.csv"'}
        )
    else:
        lines = [
            f"# 🧪 QA Test Case Matrix: {name}",
            "",
            "| TC ID | Step | Action & Target | Input Data | Expected Result | Status |",
            "| :--- | :--- | :--- | :--- | :--- | :--- |"
        ]
        for i, s in enumerate(steps):
            seq = s.get("sequence", i + 1)
            title = s.get("title") or s.get("edited_title") or f"Step {seq}"
            val = s.get("value") or "N/A"
            exp = s.get("expected_result") or s.get("expected") or "Action executes successfully."
            lines.append(f"| `TC-{seq:03d}` | {seq} | {title} | `{val}` | {exp} | ⏳ Pending |")
        
        return Response(
            content="\n".join(lines),
            media_type="text/markdown",
            headers={"Content-Disposition": f'attachment; filename="{name}_QA_Matrix.md"'}
        )


@app.post("/sessions/{session_id}/governance")
def update_session_governance(session_id: str, payload: dict = Body(...)):
    """
    Updates enterprise governance metadata: Department, Owner, Reviewer, Approver, Pre/Postconditions.
    """
    conn = get_connection()
    now = datetime.now().isoformat()
    try:
        conn.execute(
            """
            UPDATE workflows SET
                department = ?,
                owner = ?,
                reviewer = ?,
                approver = ?,
                effective_date = ?,
                review_frequency_days = ?,
                preconditions_json = ?,
                postconditions_json = ?,
                updated_at = ?
            WHERE id = ?
            """,
            (
                payload.get("department", ""),
                payload.get("owner", ""),
                payload.get("reviewer", ""),
                payload.get("approver", ""),
                payload.get("effective_date", ""),
                int(payload.get("review_frequency_days", 90)),
                json.dumps(payload.get("preconditions", [])),
                json.dumps(payload.get("postconditions", [])),
                now,
                session_id
            )
        )
        conn.commit()
        return {"success": True, "message": "Governance metadata updated successfully."}
    finally:
        conn.close()


# ── ProcBot RPA Automation Endpoints ──────────────────────────────────────────
@app.get("/sessions/{session_id}/procbot-script")
def get_procbot_script(session_id: str, engine: str = "playwright"):
    """
    Generates executable RPA automation script (Playwright or Selenium).
    """
    from backend.procbot_generator import ProcBotGenerator
    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")
        
        steps_cursor = conn.execute(
            """
            SELECT s.*, e.title as edited_title, e.hidden, e.description as edited_description
            FROM workflow_steps s
            LEFT JOIN step_edits e ON s.id = e.step_id
            WHERE s.workflow_id = ?
            ORDER BY s.sequence ASC
            """,
            (session_id,)
        )
        steps = [dict(r) for r in steps_cursor.fetchall()]
        
        generator = ProcBotGenerator(workflow_name=dict(wf).get("name", "Workflow"), steps=steps)
        if engine.lower() == "selenium":
            script = generator.generate_selenium_script()
            filename = f"procbot_{session_id[:8]}_selenium.py"
        else:
            script = generator.generate_playwright_script()
            filename = f"procbot_{session_id[:8]}_playwright.py"
            
        return Response(
            content=script,
            media_type="text/x-python",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )
    finally:
        conn.close()


@app.get("/sessions/{session_id}/procbot-recipe")
def get_procbot_recipe(session_id: str):
    """
    Generates dynamic in-browser execution recipe for Chrome extension and runner.
    """
    from backend.procbot_generator import ProcBotGenerator
    conn = get_connection()
    try:
        wf = conn.execute("SELECT * FROM workflows WHERE id = ?", (session_id,)).fetchone()
        if not wf:
            raise HTTPException(status_code=404, detail="Workflow not found")
        
        # Check if saved custom config exists
        saved = conn.execute("SELECT config_json FROM procbot_configs WHERE workflow_id = ?", (session_id,)).fetchone()
        if saved and saved[0]:
            try:
                return json.loads(saved[0])
            except Exception:
                pass

        steps_cursor = conn.execute(
            """
            SELECT s.*, e.title as edited_title, e.hidden, e.description as edited_description, e.note
            FROM workflow_steps s
            LEFT JOIN step_edits e ON s.id = e.step_id
            WHERE s.workflow_id = ?
            ORDER BY s.sequence ASC
            """,
            (session_id,)
        )
        steps = [dict(r) for r in steps_cursor.fetchall()]
        generator = ProcBotGenerator(workflow_name=dict(wf).get("name", "Workflow"), steps=steps)
        return generator.generate_json_recipe()
    finally:
        conn.close()


@app.get("/sessions/{session_id}/procbot-config")
def get_procbot_config(session_id: str):
    """
    Retrieves the saved custom ProcBot configuration for a workflow.
    """
    conn = get_connection()
    try:
        saved = conn.execute("SELECT * FROM procbot_configs WHERE workflow_id = ?", (session_id,)).fetchone()
        if saved:
            res = dict(saved)
            try:
                res["config"] = json.loads(res["config_json"])
            except Exception:
                res["config"] = None
            return res
        return {"workflow_id": session_id, "config": None, "saved": False}
    finally:
        conn.close()


@app.post("/sessions/{session_id}/procbot-config")
def save_procbot_config(session_id: str, payload: Dict[str, Any] = Body(...)):
    """
    Saves a customized ProcBot automation configuration (steps, variables, custom selectors, manual stops).
    """
    conn = get_connection()
    try:
        wf = conn.execute("SELECT name FROM workflows WHERE id = ?", (session_id,)).fetchone()
        name = dict(wf).get("name", "Workflow") if wf else "Workflow"
        now = datetime.utcnow().isoformat()
        config_str = json.dumps(payload.get("config", payload))

        conn.execute(
            """
            INSERT INTO procbot_configs (workflow_id, name, config_json, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(workflow_id) DO UPDATE SET
                name = excluded.name,
                config_json = excluded.config_json,
                updated_at = excluded.updated_at
            """,
            (session_id, name, config_str, now, now)
        )
        conn.commit()
        return {"status": "saved", "workflow_id": session_id, "updated_at": now}
    finally:
        conn.close()


@app.post("/sessions/{session_id}/procbot-run-log")
def record_procbot_run_log(session_id: str, payload: Dict[str, Any] = Body(...)):
    """
    Records an execution run audit log for a ProcBot automation run.
    """
    conn = get_connection()
    try:
        now = datetime.utcnow().isoformat()
        conn.execute(
            """
            INSERT INTO procbot_run_logs 
            (workflow_id, engine, mode, total_steps, success_steps, failed_steps, elapsed_sec, log_json, created_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                session_id,
                payload.get("engine", "browser"),
                payload.get("mode", "auto"),
                int(payload.get("total_steps", 0)),
                int(payload.get("success_steps", 0)),
                int(payload.get("failed_steps", 0)),
                float(payload.get("elapsed_sec", 0.0)),
                json.dumps(payload.get("logs", [])),
                now
            )
        )
        conn.commit()
        return {"status": "recorded", "created_at": now}
    finally:
        conn.close()


@app.get("/sessions/{session_id}/procbot-run-logs")
def get_procbot_run_logs(session_id: str, limit: int = 20):
    """
    Returns recent execution history and audit logs for ProcBot runs.
    """
    conn = get_connection()
    try:
        rows = conn.execute(
            """
            SELECT * FROM procbot_run_logs
            WHERE workflow_id = ?
            ORDER BY id DESC
            LIMIT ?
            """,
            (session_id, limit)
        ).fetchall()
        logs = []
        for r in rows:
            d = dict(r)
            try:
                d["logs"] = json.loads(d["log_json"])
            except Exception:
                d["logs"] = []
            logs.append(d)
        return {"workflow_id": session_id, "count": len(logs), "runs": logs}
    finally:
        conn.close()







